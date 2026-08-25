"""Impact Analytics brand tokens and Google Slides safe PPTX primitives.

Every value here comes from the IA Brand Guide. Nothing in this module invents a
colour, a typeface or a spacing rule. The drawing helpers below deliberately
restrict themselves to the subset of OOXML that Google Slides renders faithfully,
so a deck produced by this package looks the same in PowerPoint, Keynote and
Google Slides.
"""

from __future__ import annotations

import os
from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ─── Palette ────────────────────────────────────────────────────────────────────
# Primary palette
IMPACT_BLUE = RGBColor(0x26, 0x4C, 0xD7)
OFF_WHITE = RGBColor(0xF4, 0xF4, 0xF6)
BLACK = RGBColor(0x1C, 0x1B, 0x1B)

# Secondary palette
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_ORANGE = RGBColor(0xFF, 0x6F, 0x1C)
GRAY_1 = RGBColor(0xE4, 0xE4, 0xE4)
GRAY_2 = RGBColor(0xD7, 0xD6, 0xD2)
GRAY_3 = RGBColor(0xB3, 0xB2, 0xAD)

# Use-case palette. One solution colour per composition, never two.
SOLUTION_COLORS = {
    'inventory_replenishment': RGBColor(0xE3, 0xF5, 0x76),   # Yellow
    'merchandising': RGBColor(0xBE, 0xA8, 0xEF),             # Purple
    'pricing_promotions': RGBColor(0x3D, 0xD4, 0x99),        # Green
    'data_intelligence': RGBColor(0xB3, 0xC9, 0xF7),         # Blue
}

SOLUTION_LABELS = {
    'inventory_replenishment': 'Inventory & Replenishment',
    'merchandising': 'Merchandising',
    'pricing_promotions': 'Pricing & Promotions',
    'data_intelligence': 'Data & Intelligence',
}

# IA products mapped to their parent solution. The solution decides the colour.
PRODUCT_SOLUTIONS = {
    'ItemSmart': 'merchandising',
    'PlanSmart': 'merchandising',
    'AssortSmart': 'merchandising',
    'SizeSmart': 'merchandising',
    'VisualSmart': 'merchandising',
    'StoreSmart': 'merchandising',
    'InventorySmart': 'inventory_replenishment',
    'ForecastSmart': 'inventory_replenishment',
    'SourceSmart': 'inventory_replenishment',
    'SpaceSmart': 'inventory_replenishment',
    'PriceSmart': 'pricing_promotions',
    'PromoSmart': 'pricing_promotions',
    'MarkSmart': 'pricing_promotions',
    'BaseSmart': 'pricing_promotions',
    'TradeSmart': 'pricing_promotions',
    'AttributeSmart': 'data_intelligence',
    'CortexEye': 'data_intelligence',
    'MondaySmart': 'data_intelligence',
    'TestSmart': 'data_intelligence',
}

DEFAULT_SOLUTION = 'data_intelligence'


def solution_for_product(product: str) -> str:
    """Return the solution key that owns a product name."""
    if not product:
        return DEFAULT_SOLUTION
    for name, solution in PRODUCT_SOLUTIONS.items():
        if name.lower() == product.strip().lower():
            return solution
    return DEFAULT_SOLUTION


def solution_color(solution: str) -> RGBColor:
    return SOLUTION_COLORS.get(solution, SOLUTION_COLORS[DEFAULT_SOLUTION])


def readable_on(color: RGBColor) -> RGBColor:
    """Pick Black or White text for a background, using relative luminance."""
    r, g, b = color[0] / 255.0, color[1] / 255.0, color[2] / 255.0

    def channel(c):
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    luminance = 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)
    return BLACK if luminance > 0.45 else WHITE


def mix(color: RGBColor, other: RGBColor, weight: float) -> RGBColor:
    """Blend two palette colours. Used only for tints toward White or Off-White."""
    weight = max(0.0, min(1.0, weight))
    return RGBColor(
        int(round(color[0] * (1 - weight) + other[0] * weight)),
        int(round(color[1] * (1 - weight) + other[1] * weight)),
        int(round(color[2] * (1 - weight) + other[2] * weight)),
    )


# ─── Typography ─────────────────────────────────────────────────────────────────
# ABC Otto is licensed and rarely installed, so the brand guide specifies Inter
# Tight for PPTX headlines. Inter Tight and Spectral both ship in the Google
# Slides font list, so either choice survives the round trip.
HEADING_FONT = 'Inter Tight'
BODY_FONT = 'Inter Tight'
SERIF_HEADING_FONT = 'Spectral'

TYPE_SCALE = {
    'display': Pt(40),
    'h1': Pt(32),
    'h2': Pt(24),
    'h3': Pt(18),
    'h4': Pt(15),
    'h5': Pt(13),
    'body': Pt(12),
    'body_small': Pt(10.5),
    'caption': Pt(9),
    'micro': Pt(8),
    'stat': Pt(34),
}

# ─── Geometry ───────────────────────────────────────────────────────────────────
SLIDE_WIDTH = Emu(12192000)   # exactly 13.333in, the 16:9 canvas both renderers use
SLIDE_HEIGHT = Emu(6858000)   # exactly 7.5in
MARGIN = Inches(0.5)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN
HEADER_HEIGHT = Inches(1.05)
FOOTER_TOP = SLIDE_HEIGHT - Inches(0.42)
BODY_TOP = HEADER_HEIGHT + Inches(0.22)
BODY_HEIGHT = FOOTER_TOP - BODY_TOP - Inches(0.12)
CORNER_RADIUS = 6500  # prstGeom adjustment, roughly 0.1in on a 1in tall card

ASSET_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'static', 'images')
LOGO_PATH = os.path.normpath(os.path.join(ASSET_DIR, 'ia_logo.png'))
LOGO_WHITE_PATH = os.path.normpath(os.path.join(ASSET_DIR, 'ia_logo_white.png'))
LOGO_ASPECT = 219 / 668.0


@dataclass(frozen=True)
class Theme:
    """Resolved colour set for one battlecard. One solution colour, never two."""

    solution: str
    accent: RGBColor
    accent_text: RGBColor
    page_bg: RGBColor
    card_bg: RGBColor
    ink: RGBColor
    muted: RGBColor
    rule: RGBColor
    heading_font: str
    body_font: str

    @classmethod
    def for_solution(cls, solution: str, serif_headings: bool = False) -> 'Theme':
        accent = solution_color(solution)
        return cls(
            solution=solution,
            accent=accent,
            accent_text=readable_on(accent),
            page_bg=OFF_WHITE,
            card_bg=WHITE,
            ink=BLACK,
            muted=GRAY_3,
            rule=GRAY_1,
            heading_font=SERIF_HEADING_FONT if serif_headings else HEADING_FONT,
            body_font=BODY_FONT,
        )


# ─── Logo handling ──────────────────────────────────────────────────────────────

def ensure_white_logo() -> str:
    """Create the white logo variant the brand guide requires on blue fills.

    The source logo is transparent with an Impact Blue mark and a Black wordmark.
    Recolouring the opaque pixels to pure white keeps the shape, the spacing and
    the full opacity the guide demands.
    """
    if os.path.exists(LOGO_WHITE_PATH):
        return LOGO_WHITE_PATH
    if not os.path.exists(LOGO_PATH):
        return ''
    try:
        from PIL import Image

        src = Image.open(LOGO_PATH).convert('RGBA')
        alpha = src.getchannel('A')
        white = Image.new('RGBA', src.size, (255, 255, 255, 0))
        white.putalpha(alpha)
        os.makedirs(os.path.dirname(LOGO_WHITE_PATH), exist_ok=True)
        white.save(LOGO_WHITE_PATH, 'PNG')
        return LOGO_WHITE_PATH
    except Exception:
        return ''


def logo_for_background(dark: bool) -> str:
    """Return the logo file that belongs on this background."""
    if dark:
        path = ensure_white_logo()
        if path:
            return path
    return LOGO_PATH if os.path.exists(LOGO_PATH) else ''


# ─── Google Slides safe drawing primitives ──────────────────────────────────────

def set_run_font(run, size=None, bold=False, color=None, font=None, italic=False):
    """Set a run's typeface on every script slot.

    Google Slides reads the latin typeface and ignores theme inheritance in some
    import paths, so every run states its own font rather than relying on the
    master.
    """
    font_name = font or BODY_FONT
    run.font.name = font_name
    if size is not None:
        run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        existing = rPr.find(qn(tag))
        if existing is None:
            existing = rPr.makeelement(qn(tag), {})
            rPr.append(existing)
        existing.set('typeface', font_name)
    return run


def prepare_text_frame(shape, margin=Inches(0.14), anchor=MSO_ANCHOR.TOP):
    """Normalise a text frame so wrapping is identical across renderers."""
    tf = shape.text_frame
    tf.word_wrap = True
    # Google Slides ignores normAutofit font scaling, so this package sizes text
    # itself and turns autofit off everywhere.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = anchor
    return tf


def write_paragraph(tf, text, size, *, bold=False, color=BLACK, font=None,
                    align=PP_ALIGN.LEFT, space_before=0, space_after=4,
                    line_spacing=1.15, bullet=None, indent_level=0, first=False,
                    italic=False):
    """Append (or reuse) a paragraph and style it explicitly."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if bullet:
        apply_bullet(p, bullet, color if bullet_inherits_color(bullet) else color, indent_level)
    else:
        clear_bullet(p)
    run = p.add_run()
    run.text = text
    set_run_font(run, size, bold, color, font, italic)
    return p


def bullet_inherits_color(_bullet):
    return True


def apply_bullet(paragraph, char, color, indent_level=0):
    """Attach a character bullet with a hanging indent.

    buChar plus explicit marL/indent renders identically in PowerPoint and
    Google Slides, unlike list styles inherited from a master.
    """
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ('a:buNone', 'a:buAutoNum', 'a:buChar', 'a:buFont', 'a:buClr'):
        for node in pPr.findall(qn(tag)):
            pPr.remove(node)
    marl = Inches(0.16) + Inches(0.16) * indent_level
    pPr.set('marL', str(int(marl)))
    pPr.set('indent', str(int(-Inches(0.16))))
    buClr = pPr.makeelement(qn('a:buClr'), {})
    srgb = pPr.makeelement(qn('a:srgbClr'), {'val': str(color)})
    buClr.append(srgb)
    buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    buChar = pPr.makeelement(qn('a:buChar'), {'char': char})
    for node in (buClr, buFont, buChar):
        pPr.append(node)


def clear_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ('a:buAutoNum', 'a:buChar', 'a:buFont', 'a:buClr', 'a:buNone'):
        for node in pPr.findall(qn(tag)):
            pPr.remove(node)
    pPr.append(pPr.makeelement(qn('a:buNone'), {}))


def add_rect(slide, left, top, width, height, fill=None, line=None,
             line_width=Pt(1), rounded=False, radius=CORNER_RADIUS):
    """Add a rectangle or rounded rectangle with an explicit fill and line."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, int(left), int(top), int(width), int(height))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
    if rounded:
        set_corner_radius(shape, radius)
    strip_shape_effects(shape)
    prepare_text_frame(shape)
    return shape


def set_corner_radius(shape, value):
    """Write the rounded rectangle adjustment value that Google Slides reads."""
    prstGeom = shape._element.spPr.find(qn('a:prstGeom'))
    if prstGeom is None:
        return
    avLst = prstGeom.find(qn('a:avLst'))
    if avLst is None:
        avLst = prstGeom.makeelement(qn('a:avLst'), {})
        prstGeom.append(avLst)
    for node in list(avLst):
        avLst.remove(node)
    gd = avLst.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': 'val %d' % int(value)})
    avLst.append(gd)


def strip_shape_effects(shape):
    """Remove theme driven effects. Google Slides drops most of them anyway."""
    spPr = shape._element.spPr
    for tag in ('a:effectLst', 'a:effectRef', 'a:scene3d', 'a:sp3d'):
        for node in spPr.findall(qn(tag)):
            spPr.remove(node)
    style = shape._element.find(qn('p:style'))
    if style is not None:
        shape._element.remove(style)


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP, margin=Inches(0.0)):
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    prepare_text_frame(box, margin=margin, anchor=anchor)
    return box


def add_picture(slide, path, left, top, width=None, height=None):
    if not path or not os.path.exists(path):
        return None
    return slide.shapes.add_picture(path, int(left), int(top),
                                    None if width is None else int(width),
                                    None if height is None else int(height))


# ─── Text measurement ───────────────────────────────────────────────────────────
# Google Slides does not honour PowerPoint's shrink-on-overflow, so the builder
# measures text and picks a size that fits before the file is written.

_AVG_CHAR_RATIO = 0.50   # mean advance width of Inter Tight, relative to em
_WIDE_CHAR_RATIO = 0.58  # used for all caps labels
_LINE_HEIGHT = 1.22


def estimate_lines(text: str, width_emu: int, size_pt: float, caps: bool = False) -> int:
    """Estimate how many wrapped lines a string needs inside a box."""
    if not text:
        return 0
    ratio = _WIDE_CHAR_RATIO if caps else _AVG_CHAR_RATIO
    char_width_emu = size_pt * ratio * 12700
    if char_width_emu <= 0:
        return 1
    per_line = max(1, int(width_emu / char_width_emu))
    lines = 0
    for hard_line in text.split('\n'):
        words = hard_line.split()
        if not words:
            lines += 1
            continue
        current = 0
        line_count = 1
        for word in words:
            need = len(word) if current == 0 else len(word) + 1
            if current + need > per_line and current > 0:
                line_count += 1
                current = len(word)
            else:
                current += need
        lines += line_count
    return max(1, lines)


def text_height(text: str, width_emu: int, size_pt: float, caps: bool = False,
                extra_lines: int = 0) -> int:
    lines = estimate_lines(text, width_emu, size_pt, caps) + extra_lines
    return int(lines * size_pt * _LINE_HEIGHT * 12700)


def fit_size(text: str, width_emu: int, height_emu: int, max_pt: float,
             min_pt: float = 7.5, caps: bool = False, step: float = 0.5) -> Pt:
    """Return the largest size at or below max_pt whose wrapped text still fits."""
    size = max_pt
    while size > min_pt:
        if text_height(text, width_emu, size, caps) <= height_emu:
            break
        size -= step
    return Pt(max(min_pt, round(size, 1)))


def fit_block_size(lines, width_emu: int, height_emu: int, max_pt: float,
                   min_pt: float = 7.5, gap_pt: float = 3.0, step: float = 0.5) -> Pt:
    """Fit a list of paragraphs, allowing for the gap between each one."""
    size = max_pt
    lines = [line for line in lines if line]
    if not lines:
        return Pt(max_pt)
    while size > min_pt:
        total = sum(text_height(line, width_emu, size) for line in lines)
        total += int(gap_pt * 12700) * max(0, len(lines) - 1)
        if total <= height_emu:
            break
        size -= step
    return Pt(max(min_pt, round(size, 1)))
