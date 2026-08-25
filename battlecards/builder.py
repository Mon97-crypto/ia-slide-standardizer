"""Slide builders for the Impact Analytics competitive battlecard.

Every slide is drawn on a blank layout with explicit geometry, explicit fills and
explicit run level fonts. Nothing inherits from a theme, nothing relies on
PowerPoint autofit, and nothing uses an effect Google Slides drops on import. The
result opens the same way in PowerPoint, Keynote and Google Slides.
"""

from __future__ import annotations

import math

from lxml import etree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from . import brand
from .brand import (ACCENT_ORANGE, BLACK, GRAY_1, GRAY_2, GRAY_3, IMPACT_BLUE,
                    OFF_WHITE, WHITE, SOLUTION_LABELS, Theme, add_picture,
                    add_rect, add_textbox, fit_block_size, fit_size,
                    logo_for_background, mix, prepare_text_frame, set_run_font,
                    text_height, write_paragraph)
from .patterns import grid_overlay
from .schema import RATING_LABELS, SECTION_LABELS

SLIDE_W = Emu(12192000)   # exactly 13.333in, the 16:9 canvas Google Slides uses
SLIDE_H = Emu(6858000)    # exactly 7.5in
MARGIN = Inches(0.5)
CONTENT_W = SLIDE_W - 2 * MARGIN

KICKER_TOP = Inches(0.34)
TITLE_TOP = Inches(0.58)
HEADER_RULE_Y = Inches(1.3)
BODY_TOP = Inches(1.55)
BODY_BOTTOM = Inches(6.82)
BODY_H = BODY_BOTTOM - BODY_TOP
FOOTER_RULE_Y = Inches(6.95)
FOOTER_TEXT_Y = Inches(7.0)

GUTTER = Inches(0.22)
PAD = Inches(0.22)

RATING_COLORS = {
    'strong': IMPACT_BLUE,
    'partial': BLACK,
    'none': ACCENT_ORANGE,   # the only orange in the deck, which keeps it an accent
    'unknown': GRAY_3,
}


def chunk(items, size):
    size = max(1, size)
    return [items[i:i + size] for i in range(0, len(items), size)] or [[]]


class BattlecardDeck:
    """Turns a normalised battlecard into a Presentation."""

    def __init__(self, card: dict):
        self.card = card
        self.meta = card['meta']
        options = card.get('options', {})
        solution = self.meta.get('solution') or brand.solution_for_product(self.meta.get('ia_product', ''))
        self.solution = solution if solution in SOLUTION_LABELS else brand.DEFAULT_SOLUTION
        self.theme = Theme.for_solution(self.solution, options.get('serif_headings', False))
        self.include_notes = options.get('include_notes', True)
        self.sections = options.get('sections') or list(SECTION_LABELS)
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._apply_theme_fonts()
        self.page = 0

    # ── presentation level ─────────────────────────────────────────────────────

    def _apply_theme_fonts(self):
        """Name the brand fonts in the theme as well as on every run.

        Runs already carry their own typeface, so rendering never depends on this.
        Setting the theme too means a user who edits the deck in Google Slides or
        PowerPoint gets Inter Tight on any new text box.
        """
        theme_rel = ('http://schemas.openxmlformats.org/officeDocument/2006/'
                     'relationships/theme')
        try:
            part = self.prs.slide_masters[0].part.part_related_by(theme_rel)
            root = etree.fromstring(part.blob)
        except Exception:
            return
        font_scheme = root.find('.//' + qn('a:fontScheme'))
        if font_scheme is None:
            return
        for tag, name in (('a:majorFont', self.theme.heading_font),
                          ('a:minorFont', self.theme.body_font)):
            node = font_scheme.find(qn(tag))
            if node is None:
                continue
            latin = node.find(qn('a:latin'))
            if latin is not None:
                latin.set('typeface', name)
        try:
            part._blob = etree.tostring(root, xml_declaration=True,
                                        encoding='UTF-8', standalone=True)
        except Exception:
            return

    def build(self) -> Presentation:
        builders = {
            'cover': self.slide_cover,
            'how_to_use': self.slide_how_to_use,
            'snapshot': self.slide_snapshot,
            'positioning': self.slide_positioning,
            'strengths_weaknesses': self.slide_strengths_weaknesses,
            'why_we_win': self.slide_why_we_win,
            'comparison': self.slide_comparison,
            'objections': self.slide_objections,
            'landmines': self.slide_landmines,
            'discovery': self.slide_discovery,
            'proof_points': self.slide_proof_points,
            'talk_track': self.slide_talk_track,
            'dos_donts': self.slide_dos_donts,
            'pricing': self.slide_pricing,
            'next_steps': self.slide_next_steps,
            'one_pager': self.slide_one_pager,
        }
        for name in self.sections:
            builder = builders.get(name)
            if builder:
                builder()
        return self.prs

    # ── slide chrome ───────────────────────────────────────────────────────────

    def _slide(self, dark=False, patterned=False):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = IMPACT_BLUE if dark else OFF_WHITE
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=bg)
        if patterned:
            overlay = grid_overlay()
            if overlay:
                add_picture(slide, overlay, 0, 0, SLIDE_W, SLIDE_H)
        return slide

    def _header(self, slide, title, kicker=None, dark=False):
        ink = WHITE if dark else BLACK
        kicker_color = WHITE if dark else IMPACT_BLUE
        kicker_text = kicker if kicker is not None else self._default_kicker()

        box = add_textbox(slide, MARGIN, KICKER_TOP, CONTENT_W - Inches(2.0), Inches(0.24))
        write_paragraph(box.text_frame, kicker_text.upper(), Pt(9), bold=True,
                        color=kicker_color, font=self.theme.body_font,
                        space_after=0, line_spacing=1.0, first=True)

        title_box = add_textbox(slide, MARGIN, TITLE_TOP, CONTENT_W - Inches(2.0), Inches(0.66))
        size = fit_size(title, int(CONTENT_W - Inches(2.0)), int(Inches(0.66)), 24, 15)
        write_paragraph(title_box.text_frame, title, size, bold=True, color=ink,
                        font=self.theme.heading_font, space_after=0,
                        line_spacing=1.05, first=True)

        add_rect(slide, MARGIN, HEADER_RULE_Y, Inches(1.1), Inches(0.045), fill=self.theme.accent)
        add_rect(slide, MARGIN + Inches(1.1), HEADER_RULE_Y + Inches(0.018),
                 CONTENT_W - Inches(1.1), Inches(0.01),
                 fill=WHITE if dark else GRAY_1)
        self._logo(slide, dark)

    def _logo(self, slide, dark=False, height=Inches(0.3)):
        path = logo_for_background(dark)
        if not path:
            return
        width = int(height / brand.LOGO_ASPECT)
        add_picture(slide, path, SLIDE_W - MARGIN - width, KICKER_TOP - Inches(0.02),
                    width, height)

    def _footer(self, slide, dark=False):
        self.page += 1
        ink = mix(WHITE, IMPACT_BLUE, 0.35) if dark else GRAY_3
        add_rect(slide, MARGIN, FOOTER_RULE_Y, CONTENT_W, Inches(0.01),
                 fill=mix(WHITE, IMPACT_BLUE, 0.6) if dark else GRAY_1)
        left = add_textbox(slide, MARGIN, FOOTER_TEXT_Y, CONTENT_W - Inches(1.2), Inches(0.28))
        label = '%s vs Impact Analytics  ·  %s  ·  %s' % (
            self.meta['competitor'], self.meta['confidentiality'], self.meta['version'])
        write_paragraph(left.text_frame, label, Pt(8), color=ink,
                        font=self.theme.body_font, space_after=0, line_spacing=1.0, first=True)
        right = add_textbox(slide, SLIDE_W - MARGIN - Inches(1.0), FOOTER_TEXT_Y,
                            Inches(1.0), Inches(0.28))
        write_paragraph(right.text_frame, str(self.page), Pt(8), color=ink,
                        font=self.theme.body_font, align=PP_ALIGN.RIGHT,
                        space_after=0, line_spacing=1.0, first=True)

    def _default_kicker(self):
        return '%s battlecard  ·  %s' % (
            SOLUTION_LABELS.get(self.solution, ''), self.meta['competitor'])

    def _notes(self, slide, text):
        if not self.include_notes or not text:
            return
        slide.notes_slide.notes_text_frame.text = text

    def _page(self, title, kicker=None, dark=False, patterned=False):
        slide = self._slide(dark=dark, patterned=patterned)
        self._header(slide, title, kicker, dark=dark)
        self._footer(slide, dark=dark)
        return slide

    # ── reusable blocks ────────────────────────────────────────────────────────

    def _panel(self, slide, left, top, width, height, fill=None, border=GRAY_1,
               accent=None, accent_height=Inches(0.05)):
        panel = add_rect(slide, left, top, width, height,
                         fill=fill or WHITE, line=border, line_width=Pt(0.75), rounded=True)
        if accent is not None:
            add_rect(slide, left, top, width, accent_height, fill=accent, rounded=True,
                     radius=1500)
        return panel

    def _panel_label(self, slide, left, top, width, text, color=None, size=Pt(10)):
        box = add_textbox(slide, left, top, width, Inches(0.26))
        write_paragraph(box.text_frame, text.upper(), size, bold=True,
                        color=color or IMPACT_BLUE, font=self.theme.body_font,
                        space_after=0, line_spacing=1.0, first=True)
        return box

    def _bullets(self, slide, left, top, width, height, items, *, max_pt=12.0,
                 min_pt=8.0, color=BLACK, bullet_color=None, bullet='•'):
        if not items:
            return
        inner = int(width - 2 * Inches(0.16))
        size = fit_block_size(items, inner, int(height), max_pt, min_pt, gap_pt=5.0)
        box = add_textbox(slide, left, top, width, height)
        tf = box.text_frame
        tf.margin_left = Inches(0.0)
        tf.margin_right = Inches(0.0)
        for index, item in enumerate(items):
            write_paragraph(tf, item, size, color=color, font=self.theme.body_font,
                            bullet=bullet, space_after=4, line_spacing=1.2,
                            first=index == 0)
            if bullet_color is not None:
                brand.apply_bullet(tf.paragraphs[index], bullet, bullet_color)

    def _paragraph_block(self, slide, left, top, width, height, text, *,
                         max_pt=12.0, min_pt=8.0, color=BLACK, bold=False,
                         align=PP_ALIGN.LEFT, font=None):
        if not text:
            return
        size = fit_size(text, int(width), int(height), max_pt, min_pt)
        box = add_textbox(slide, left, top, width, height)
        write_paragraph(box.text_frame, text, size, bold=bold, color=color,
                        font=font or self.theme.body_font, align=align,
                        space_after=0, line_spacing=1.25, first=True)

    def _stat(self, slide, left, top, width, value, label, detail='', source=''):
        height = Inches(1.9)
        self._panel(slide, left, top, width, height, accent=self.theme.accent)
        inner_w = width - 2 * PAD
        value_size = fit_size(value or '', int(inner_w), int(Inches(0.75)), 34, 16)
        box = add_textbox(slide, left + PAD, top + Inches(0.24), inner_w, Inches(0.72))
        write_paragraph(box.text_frame, value or 'Add stat', value_size, bold=True,
                        color=IMPACT_BLUE, font=self.theme.heading_font,
                        space_after=0, line_spacing=1.0, first=True)
        label_box = add_textbox(slide, left + PAD, top + Inches(0.96), inner_w, Inches(0.3))
        write_paragraph(label_box.text_frame, label or '', Pt(11), bold=True, color=BLACK,
                        font=self.theme.body_font, space_after=0, line_spacing=1.1, first=True)
        if detail:
            self._paragraph_block(slide, left + PAD, top + Inches(1.26), inner_w,
                                  Inches(0.44), detail, max_pt=9.5, min_pt=7.5, color=GRAY_3)
        if source:
            self._link(slide, left + PAD, top + height - Inches(0.34), inner_w,
                       'Source', source)

    def _link(self, slide, left, top, width, label, url):
        box = add_textbox(slide, left, top, width, Inches(0.24))
        p = box.text_frame.paragraphs[0]
        p.line_spacing = 1.0
        run = p.add_run()
        run.text = label
        set_run_font(run, Pt(8.5), False, IMPACT_BLUE, self.theme.body_font)
        if url.lower().startswith('http'):
            run.hyperlink.address = url
        else:
            run.text = url[:90]
        return box

    def _table(self, slide, left, top, width, headers, rows, col_ratios,
               row_height=Inches(0.44), header_height=Inches(0.4), font_pt=10.0):
        """Draw a table with explicit fills, borders and fonts.

        The theme table style is removed so Google Slides never substitutes its
        own banding.
        """
        total = len(rows) + 1
        frame = slide.shapes.add_table(total, len(headers), int(left), int(top),
                                       int(width), int(header_height + row_height * len(rows)))
        table = frame.table
        tblPr = table._tbl.find(qn('a:tblPr'))
        if tblPr is not None:
            for node in tblPr.findall(qn('a:tableStyleId')):
                tblPr.remove(node)
            tblPr.set('firstRow', '0')
            tblPr.set('bandRow', '0')

        ratio_total = float(sum(col_ratios))
        for index, ratio in enumerate(col_ratios):
            table.columns[index].width = Emu(int(width * ratio / ratio_total))
        table.rows[0].height = int(header_height)
        for index in range(1, total):
            table.rows[index].height = int(row_height)

        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            self._cell(cell, header, size=Pt(9.5), bold=True, color=WHITE,
                       fill=IMPACT_BLUE, align=PP_ALIGN.LEFT if col == 0 else PP_ALIGN.CENTER)

        for r, row in enumerate(rows, start=1):
            band = WHITE if r % 2 else OFF_WHITE
            for c, value in enumerate(row):
                text, color, bold, align = self._unpack_cell(value, c)
                cell = table.cell(r, c)
                self._cell(cell, text, size=Pt(font_pt), bold=bold, color=color,
                           fill=band, align=align)
        return table

    @staticmethod
    def _unpack_cell(value, column_index):
        if isinstance(value, dict):
            return (value.get('text', ''), value.get('color', BLACK),
                    value.get('bold', False),
                    value.get('align', PP_ALIGN.LEFT if column_index == 0 else PP_ALIGN.CENTER))
        return (str(value), BLACK, False,
                PP_ALIGN.LEFT if column_index == 0 else PP_ALIGN.CENTER)

    def _cell(self, cell, text, size, bold, color, fill, align):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = 1.12
        run = p.add_run()
        run.text = text
        set_run_font(run, size, bold, color, self.theme.body_font)
        self._cell_borders(cell, GRAY_1)

    @staticmethod
    def _cell_borders(cell, color, width=Pt(0.75)):
        tcPr = cell._tc.get_or_add_tcPr()
        order = ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']
        for tag in reversed(order):
            for node in tcPr.findall(qn(tag)):
                tcPr.remove(node)
            line = tcPr.makeelement(qn(tag), {'w': str(int(width)), 'cap': 'flat',
                                              'cmpd': 'sng', 'algn': 'ctr'})
            solid = line.makeelement(qn('a:solidFill'), {})
            srgb = line.makeelement(qn('a:srgbClr'), {'val': str(color)})
            solid.append(srgb)
            line.append(solid)
            tcPr.insert(0, line)

    # ── sections ───────────────────────────────────────────────────────────────

    def slide_cover(self):
        slide = self._slide(dark=True, patterned=True)
        self._logo(slide, dark=True, height=Inches(0.42))

        add_rect(slide, MARGIN, Inches(2.05), Inches(1.4), Inches(0.05), fill=self.theme.accent)

        kicker = 'Competitive battlecard  ·  %s' % SOLUTION_LABELS.get(self.solution, '')
        box = add_textbox(slide, MARGIN, Inches(1.68), CONTENT_W, Inches(0.3))
        write_paragraph(box.text_frame, kicker.upper(), Pt(10), bold=True, color=WHITE,
                        font=self.theme.body_font, space_after=0, line_spacing=1.0, first=True)

        title = self.meta['competitor']
        title_size = fit_size(title, int(CONTENT_W), int(Inches(1.5)), 54, 26)
        title_box = add_textbox(slide, MARGIN, Inches(2.35), CONTENT_W, Inches(1.5))
        write_paragraph(title_box.text_frame, title, title_size, bold=True, color=WHITE,
                        font=self.theme.heading_font, space_after=0, line_spacing=1.0, first=True)

        headline = self.meta.get('headline') or 'Know the rival. Lead with the outcome.'
        self._paragraph_block(slide, MARGIN, Inches(3.9), Inches(7.4), Inches(0.95),
                              headline, max_pt=16, min_pt=11,
                              color=mix(WHITE, IMPACT_BLUE, 0.12))

        facts = [
            ('Impact Analytics product', self.meta.get('ia_product') or 'Portfolio'),
            ('Their category', self.meta.get('competitor_category') or 'Add the category'),
            ('Prepared for', self.meta.get('audience')),
            ('Owner', self.meta.get('owner') or 'Add an owner'),
            ('Updated', self.meta.get('date')),
        ]
        top = Inches(5.05)
        col_w = (CONTENT_W - 4 * Inches(0.2)) / 5
        for index, (label, value) in enumerate(facts):
            left = MARGIN + index * (col_w + Inches(0.2))
            label_box = add_textbox(slide, left, top, col_w, Inches(0.24))
            write_paragraph(label_box.text_frame, label.upper(), Pt(8), bold=True,
                            color=mix(WHITE, IMPACT_BLUE, 0.35), font=self.theme.body_font,
                            space_after=0, line_spacing=1.0, first=True)
            self._paragraph_block(slide, left, top + Inches(0.24), col_w, Inches(0.6),
                                  value, max_pt=11, min_pt=8, color=WHITE)

        strip = add_textbox(slide, MARGIN, Inches(6.72), CONTENT_W, Inches(0.3))
        write_paragraph(strip.text_frame, self.meta['confidentiality'].upper(), Pt(9),
                        bold=True, color=mix(WHITE, IMPACT_BLUE, 0.35),
                        font=self.theme.body_font, space_after=0, line_spacing=1.0, first=True)
        self.page += 1
        self._notes(slide, 'Battlecard for %s. Owner: %s. Confirm every competitor claim '
                           'against a dated source before you use it in a deal.'
                    % (self.meta['competitor'], self.meta.get('owner') or 'unassigned'))

    def slide_how_to_use(self):
        items = self.card.get('how_to_use') or []
        if not items:
            return
        slide = self._page('How to use this battlecard')
        left_w = Inches(7.6)
        self._panel(slide, MARGIN, BODY_TOP, left_w, BODY_H, accent=self.theme.accent)
        self._panel_label(slide, MARGIN + PAD, BODY_TOP + Inches(0.3), left_w - 2 * PAD,
                          'Working rules')
        self._bullets(slide, MARGIN + PAD, BODY_TOP + Inches(0.62), left_w - 2 * PAD,
                      BODY_H - Inches(0.9), items, max_pt=13, min_pt=9,
                      bullet_color=IMPACT_BLUE)

        right_left = MARGIN + left_w + GUTTER
        right_w = CONTENT_W - left_w - GUTTER
        self._panel(slide, right_left, BODY_TOP, right_w, BODY_H, fill=IMPACT_BLUE, border=None)
        self._panel_label(slide, right_left + PAD, BODY_TOP + Inches(0.3), right_w - 2 * PAD,
                          'Win theme', color=mix(WHITE, IMPACT_BLUE, 0.35))
        self._paragraph_block(slide, right_left + PAD, BODY_TOP + Inches(0.66),
                              right_w - 2 * PAD, BODY_H - Inches(1.0),
                              self.meta.get('win_theme') or 'Name the one reason this buyer switches.',
                              max_pt=15, min_pt=10, color=WHITE)
        self._notes(slide, 'Read this slide before the call, not during it.')

    def slide_snapshot(self):
        snapshot = self.card.get('snapshot', {})
        fields = [
            ('Headquarters', snapshot.get('headquarters')),
            ('Founded', snapshot.get('founded')),
            ('Employees', snapshot.get('employees')),
            ('Ownership', snapshot.get('ownership')),
            ('Funding', snapshot.get('funding')),
            ('Target segment', snapshot.get('target_segment')),
            ('Go to market', snapshot.get('go_to_market')),
            ('Deployment', snapshot.get('deployment')),
        ]
        filled = [(label, value) for label, value in fields if value]
        moves = snapshot.get('recent_moves') or []
        customers = snapshot.get('notable_customers')
        if not filled and not moves and not customers:
            return

        slide = self._page('Competitor snapshot')
        left_w = Inches(7.6)
        self._panel(slide, MARGIN, BODY_TOP, left_w, BODY_H, accent=self.theme.accent)
        self._panel_label(slide, MARGIN + PAD, BODY_TOP + Inches(0.3), left_w - 2 * PAD,
                          'Company facts')

        grid_top = BODY_TOP + Inches(0.66)
        cell_w = (left_w - 2 * PAD - Inches(0.3)) / 2
        cell_h = Inches(0.72)
        display = filled or [(label, 'Add detail') for label, _ in fields[:6]]
        for index, (label, value) in enumerate(display[:8]):
            col, row = index % 2, index // 2
            left = MARGIN + PAD + col * (cell_w + Inches(0.3))
            top = grid_top + row * cell_h
            label_box = add_textbox(slide, left, top, cell_w, Inches(0.22))
            write_paragraph(label_box.text_frame, label.upper(), Pt(8), bold=True,
                            color=GRAY_3, font=self.theme.body_font,
                            space_after=0, line_spacing=1.0, first=True)
            self._paragraph_block(slide, left, top + Inches(0.2), cell_w, Inches(0.44),
                                  value, max_pt=11.5, min_pt=8)

        if customers:
            top = grid_top + math.ceil(len(display[:8]) / 2) * cell_h + Inches(0.06)
            available = BODY_TOP + BODY_H - top - Inches(0.2)
            if available > Inches(0.4):
                label_box = add_textbox(slide, MARGIN + PAD, top, left_w - 2 * PAD, Inches(0.22))
                write_paragraph(label_box.text_frame, 'NAMED CUSTOMERS', Pt(8), bold=True,
                                color=GRAY_3, font=self.theme.body_font,
                                space_after=0, line_spacing=1.0, first=True)
                self._paragraph_block(slide, MARGIN + PAD, top + Inches(0.2),
                                      left_w - 2 * PAD, available - Inches(0.2),
                                      customers, max_pt=11, min_pt=8)

        right_left = MARGIN + left_w + GUTTER
        right_w = CONTENT_W - left_w - GUTTER
        self._panel(slide, right_left, BODY_TOP, right_w, BODY_H, fill=WHITE)
        self._panel_label(slide, right_left + PAD, BODY_TOP + Inches(0.3),
                          right_w - 2 * PAD, 'Recent moves')
        self._bullets(slide, right_left + PAD, BODY_TOP + Inches(0.62), right_w - 2 * PAD,
                      BODY_H - Inches(0.9),
                      moves or ['Add funding, product or leadership news, with a source link.'],
                      max_pt=11.5, min_pt=8.5, bullet_color=self.theme.accent)
        self._notes(slide, 'Every fact on this slide needs a dated source. Refresh it each quarter.')

    def slide_positioning(self):
        positioning = self.card.get('positioning', {})
        if not any(positioning.values()):
            return
        slide = self._page('Positioning face off')
        col_w = (CONTENT_W - GUTTER) / 2
        panel_h = BODY_H - Inches(1.35)

        self._panel(slide, MARGIN, BODY_TOP, col_w, panel_h, fill=WHITE, accent=GRAY_2)
        self._panel_label(slide, MARGIN + PAD, BODY_TOP + Inches(0.3), col_w - 2 * PAD,
                          'They say', color=GRAY_3)
        self._paragraph_block(slide, MARGIN + PAD, BODY_TOP + Inches(0.66), col_w - 2 * PAD,
                              panel_h - Inches(0.95),
                              positioning.get('their_claim') or 'Paste their positioning line.',
                              max_pt=15, min_pt=10)

        right = MARGIN + col_w + GUTTER
        self._panel(slide, right, BODY_TOP, col_w, panel_h, fill=IMPACT_BLUE, border=None)
        self._panel_label(slide, right + PAD, BODY_TOP + Inches(0.3), col_w - 2 * PAD,
                          'We say', color=mix(WHITE, IMPACT_BLUE, 0.35))
        self._paragraph_block(slide, right + PAD, BODY_TOP + Inches(0.66), col_w - 2 * PAD,
                              panel_h - Inches(0.95),
                              positioning.get('our_claim') or 'Write the Impact Analytics line.',
                              max_pt=15, min_pt=10, color=WHITE)

        wedge_top = BODY_TOP + panel_h + Inches(0.2)
        self._panel(slide, MARGIN, wedge_top, CONTENT_W, Inches(1.15),
                    fill=WHITE, accent=self.theme.accent)
        self._panel_label(slide, MARGIN + PAD, wedge_top + Inches(0.24), CONTENT_W - 2 * PAD,
                          'The wedge')
        self._paragraph_block(slide, MARGIN + PAD, wedge_top + Inches(0.54),
                              CONTENT_W - 2 * PAD, Inches(0.5),
                              positioning.get('wedge') or 'Name the gap this buyer feels weekly.',
                              max_pt=13, min_pt=9)
        self._notes(slide, 'Say the wedge in the buyer own words. Quote them back.')

    def slide_strengths_weaknesses(self):
        strengths = self.card.get('their_strengths') or []
        weaknesses = self.card.get('their_weaknesses') or []
        if not strengths and not weaknesses:
            return
        slide = self._page('Where they win, where they fall short')
        col_w = (CONTENT_W - GUTTER) / 2
        for index, (label, items, accent) in enumerate((
                ('Where %s wins' % self.meta['competitor'], strengths, GRAY_2),
                ('Where %s falls short' % self.meta['competitor'], weaknesses, self.theme.accent))):
            left = MARGIN + index * (col_w + GUTTER)
            self._panel(slide, left, BODY_TOP, col_w, BODY_H, accent=accent)
            self._panel_label(slide, left + PAD, BODY_TOP + Inches(0.3), col_w - 2 * PAD,
                              label, color=GRAY_3 if index == 0 else IMPACT_BLUE)
            self._bullets(slide, left + PAD, BODY_TOP + Inches(0.66), col_w - 2 * PAD,
                          BODY_H - Inches(0.95),
                          items or ['Add at least three points.'],
                          max_pt=13, min_pt=9,
                          bullet_color=GRAY_3 if index == 0 else IMPACT_BLUE)
        self._notes(slide, 'An honest read of their strengths buys credibility for the rest.')

    def slide_why_we_win(self):
        advantages = self.card.get('our_advantages') or []
        if not advantages:
            return
        for group in chunk(advantages, 3):
            slide = self._page('Why Impact Analytics wins')
            count = len(group)
            col_w = (CONTENT_W - GUTTER * (count - 1)) / count
            for index, row in enumerate(group):
                left = MARGIN + index * (col_w + GUTTER)
                self._panel(slide, left, BODY_TOP, col_w, BODY_H, accent=self.theme.accent)
                number = add_textbox(slide, left + PAD, BODY_TOP + Inches(0.3),
                                     col_w - 2 * PAD, Inches(0.4))
                write_paragraph(number.text_frame, '0%d' % (index + 1), Pt(20), bold=True,
                                color=self.theme.accent, font=self.theme.heading_font,
                                space_after=0, line_spacing=1.0, first=True)
                self._paragraph_block(slide, left + PAD, BODY_TOP + Inches(0.78),
                                      col_w - 2 * PAD, Inches(0.85),
                                      row.get('title', ''), max_pt=17, min_pt=12,
                                      bold=True, font=self.theme.heading_font)
                self._paragraph_block(slide, left + PAD, BODY_TOP + Inches(1.72),
                                      col_w - 2 * PAD, BODY_H - Inches(2.9),
                                      row.get('detail', ''), max_pt=12, min_pt=8.5)
                proof = row.get('proof')
                if proof:
                    proof_top = BODY_TOP + BODY_H - Inches(1.1)
                    add_rect(slide, left + PAD, proof_top, col_w - 2 * PAD, Inches(0.01),
                             fill=GRAY_1)
                    label = add_textbox(slide, left + PAD, proof_top + Inches(0.1),
                                        col_w - 2 * PAD, Inches(0.22))
                    write_paragraph(label.text_frame, 'PROOF', Pt(8), bold=True,
                                    color=GRAY_3, font=self.theme.body_font,
                                    space_after=0, line_spacing=1.0, first=True)
                    self._paragraph_block(slide, left + PAD, proof_top + Inches(0.32),
                                          col_w - 2 * PAD, Inches(0.62), proof,
                                          max_pt=10, min_pt=7.5, color=IMPACT_BLUE)
            self._notes(slide, 'Lead with the outcome. Name the product second.')

    def slide_comparison(self):
        rows = self.card.get('comparison') or []
        if not rows:
            return
        pages = chunk(rows, 9)
        for index, group in enumerate(pages):
            title = 'Head to head'
            if len(pages) > 1:
                title = 'Head to head (%d of %d)' % (index + 1, len(pages))
            slide = self._page(title)
            headers = ['Capability', 'Impact Analytics', self.meta['competitor'], 'What to say']
            table_rows = []
            for row in group:
                table_rows.append([
                    {'text': row.get('capability', ''), 'bold': True, 'align': PP_ALIGN.LEFT},
                    self._rating_cell(row.get('ia')),
                    self._rating_cell(row.get('competitor')),
                    {'text': row.get('note', ''), 'align': PP_ALIGN.LEFT, 'color': BLACK},
                ])
            row_height = min(Inches(0.52), max(Inches(0.34),
                                               (BODY_H - Inches(0.55)) / max(1, len(group))))
            self._table(slide, MARGIN, BODY_TOP, CONTENT_W, headers, table_rows,
                        col_ratios=[0.30, 0.14, 0.14, 0.42],
                        row_height=row_height, font_pt=9.5)
            legend_top = BODY_TOP + Inches(0.42) + row_height * len(group) + Inches(0.14)
            if legend_top < FOOTER_RULE_Y - Inches(0.3):
                self._legend(slide, MARGIN, legend_top)
            self._notes(slide, 'Show this only when the buyer asks for a direct comparison. '
                               'Defend every row with evidence.')

    def _rating_cell(self, rating):
        key = rating if rating in RATING_LABELS else 'unknown'
        return {'text': RATING_LABELS[key], 'color': RATING_COLORS[key],
                'bold': key == 'strong', 'align': PP_ALIGN.CENTER}

    def _legend(self, slide, left, top):
        box = add_textbox(slide, left, top, CONTENT_W, Inches(0.26))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.line_spacing = 1.0
        for index, key in enumerate(('strong', 'partial', 'none', 'unknown')):
            run = p.add_run()
            run.text = ('   ' if index else '') + RATING_LABELS[key]
            set_run_font(run, Pt(8.5), True, RATING_COLORS[key], self.theme.body_font)
            gloss = p.add_run()
            gloss.text = {'strong': ' ships today', 'partial': ' partial coverage',
                          'none': ' not available', 'unknown': ' needs research'}[key]
            set_run_font(gloss, Pt(8.5), False, GRAY_3, self.theme.body_font)

    def slide_objections(self):
        rows = self.card.get('objections') or []
        if not rows:
            return
        for group in chunk(rows, 3):
            slide = self._page('Objection handling')
            count = len(group)
            row_h = (BODY_H - GUTTER * (count - 1)) / count
            for index, row in enumerate(group):
                top = BODY_TOP + index * (row_h + GUTTER)
                self._panel(slide, MARGIN, top, CONTENT_W, row_h, fill=WHITE)
                add_rect(slide, MARGIN, top, Inches(0.06), row_h, fill=self.theme.accent)
                col_w = (CONTENT_W - Inches(0.4)) / 3
                blocks = (
                    ('They say', row.get('objection', ''), GRAY_3, BLACK),
                    ('We say', row.get('response', ''), IMPACT_BLUE, BLACK),
                    ('Proof', row.get('proof', ''), GRAY_3, IMPACT_BLUE),
                )
                for col, (label, text, label_color, text_color) in enumerate(blocks):
                    left = MARGIN + Inches(0.28) + col * col_w
                    inner_w = col_w - Inches(0.24)
                    self._panel_label(slide, left, top + Inches(0.16), inner_w, label,
                                      color=label_color, size=Pt(8.5))
                    self._paragraph_block(slide, left, top + Inches(0.44), inner_w,
                                          row_h - Inches(0.62), text,
                                          max_pt=11.5, min_pt=8, color=text_color)
                    if col < 2:
                        add_rect(slide, left + inner_w + Inches(0.11), top + Inches(0.16),
                                 Inches(0.008), row_h - Inches(0.32), fill=GRAY_1)
            self._notes(slide, 'Answer the objection once, then return to the outcome.')

    def slide_landmines(self):
        rows = self.card.get('landmines') or []
        if not rows:
            return
        for group in chunk(rows, 3):
            slide = self._page('Landmines to set')
            count = len(group)
            col_w = (CONTENT_W - GUTTER * (count - 1)) / count
            for index, row in enumerate(group):
                left = MARGIN + index * (col_w + GUTTER)
                self._panel(slide, left, BODY_TOP, col_w, BODY_H, accent=self.theme.accent)
                self._panel_label(slide, left + PAD, BODY_TOP + Inches(0.3), col_w - 2 * PAD,
                                  'Ask this')
                self._paragraph_block(slide, left + PAD, BODY_TOP + Inches(0.62),
                                      col_w - 2 * PAD, Inches(1.5),
                                      row.get('question', ''), max_pt=14, min_pt=10,
                                      bold=True, font=self.theme.heading_font)
                add_rect(slide, left + PAD, BODY_TOP + Inches(2.2), col_w - 2 * PAD,
                         Inches(0.01), fill=GRAY_1)
                self._panel_label(slide, left + PAD, BODY_TOP + Inches(2.34),
                                  col_w - 2 * PAD, 'Why it lands', color=GRAY_3, size=Pt(8.5))
                self._paragraph_block(slide, left + PAD, BODY_TOP + Inches(2.62),
                                      col_w - 2 * PAD, Inches(1.15),
                                      row.get('why', ''), max_pt=11, min_pt=8)
                self._panel_label(slide, left + PAD, BODY_TOP + Inches(3.85),
                                  col_w - 2 * PAD, 'Listen for', color=GRAY_3, size=Pt(8.5))
                self._paragraph_block(slide, left + PAD, BODY_TOP + Inches(4.13),
                                      col_w - 2 * PAD, BODY_H - Inches(4.35),
                                      row.get('listen_for', ''), max_pt=11, min_pt=8,
                                      color=IMPACT_BLUE)
            self._notes(slide, 'Set one landmine per call. More than one sounds rehearsed.')

    def slide_discovery(self):
        rows = self.card.get('discovery') or []
        if not rows:
            return
        for group in chunk(rows, 3):
            slide = self._page('Discovery questions')
            count = len(group)
            col_w = (CONTENT_W - GUTTER * (count - 1)) / count
            for index, row in enumerate(group):
                left = MARGIN + index * (col_w + GUTTER)
                self._panel(slide, left, BODY_TOP, col_w, BODY_H, accent=self.theme.accent)
                self._paragraph_block(slide, left + PAD, BODY_TOP + Inches(0.3),
                                      col_w - 2 * PAD, Inches(0.72),
                                      row.get('theme', ''), max_pt=16, min_pt=11,
                                      bold=True, font=self.theme.heading_font)
                add_rect(slide, left + PAD, BODY_TOP + Inches(1.08), col_w - 2 * PAD,
                         Inches(0.01), fill=GRAY_1)
                self._bullets(slide, left + PAD, BODY_TOP + Inches(1.24), col_w - 2 * PAD,
                              BODY_H - Inches(1.5), row.get('questions') or [],
                              max_pt=12, min_pt=8.5, bullet_color=IMPACT_BLUE)
            self._notes(slide, 'Run discovery before any slide goes on the screen.')

    def slide_proof_points(self):
        rows = self.card.get('proof_points') or []
        if not rows:
            return
        for group in chunk(rows, 4):
            slide = self._page('Proof points')
            count = len(group)
            col_w = (CONTENT_W - GUTTER * (count - 1)) / count
            for index, row in enumerate(group):
                left = MARGIN + index * (col_w + GUTTER)
                self._stat(slide, left, BODY_TOP, col_w, row.get('stat', ''),
                           row.get('label', ''), row.get('detail', ''), row.get('source', ''))
            note_top = BODY_TOP + Inches(2.1)
            self._panel(slide, MARGIN, note_top, CONTENT_W, BODY_H - Inches(2.1),
                        fill=WHITE, accent=self.theme.accent)
            self._panel_label(slide, MARGIN + PAD, note_top + Inches(0.24),
                              CONTENT_W - 2 * PAD, 'How to use these numbers')
            self._bullets(slide, MARGIN + PAD, note_top + Inches(0.58), CONTENT_W - 2 * PAD,
                          BODY_H - Inches(2.9), [
                              'Quote the number, then name the customer situation behind it.',
                              'Carry the source link. A stat without a source loses the room.',
                              'Use results from 2025 or later. Retire anything older.',
                          ], max_pt=12, min_pt=9, bullet_color=IMPACT_BLUE)
            self._notes(slide, 'Never quote a result you cannot source on request.')

    def slide_talk_track(self):
        track = self.card.get('talk_track', {})
        blocks = [
            ('Positioning statement', track.get('positioning', '')),
            ('Thirty second pitch', track.get('elevator', '')),
            ('Discovery opener', track.get('discovery_open', '')),
            ('The trap question', track.get('trap', '')),
        ]
        blocks = [(label, text) for label, text in blocks if text]
        if not blocks:
            return
        slide = self._page('Talk track')
        cols = 2
        rows = math.ceil(len(blocks) / cols)
        col_w = (CONTENT_W - GUTTER) / cols
        row_h = (BODY_H - GUTTER * (rows - 1)) / rows
        for index, (label, text) in enumerate(blocks):
            col, row = index % cols, index // cols
            left = MARGIN + col * (col_w + GUTTER)
            top = BODY_TOP + row * (row_h + GUTTER)
            highlight = index == 0
            self._panel(slide, left, top, col_w, row_h,
                        fill=IMPACT_BLUE if highlight else WHITE,
                        border=None if highlight else GRAY_1,
                        accent=None if highlight else self.theme.accent)
            self._panel_label(slide, left + PAD, top + Inches(0.24), col_w - 2 * PAD, label,
                              color=mix(WHITE, IMPACT_BLUE, 0.35) if highlight else IMPACT_BLUE)
            self._paragraph_block(slide, left + PAD, top + Inches(0.58), col_w - 2 * PAD,
                                  row_h - Inches(0.82), text, max_pt=13.5, min_pt=9,
                                  color=WHITE if highlight else BLACK)
        self._notes(slide, 'Say it out loud twice before the call. Cut any sentence that does not land.')

    def slide_dos_donts(self):
        dos = self.card.get('dos') or []
        donts = self.card.get('donts') or []
        if not dos and not donts:
            return
        slide = self._page('Do and do not')
        col_w = (CONTENT_W - GUTTER) / 2
        for index, (label, items, accent, color) in enumerate((
                ('Do', dos, self.theme.accent, IMPACT_BLUE),
                ('Do not', donts, GRAY_2, GRAY_3))):
            left = MARGIN + index * (col_w + GUTTER)
            self._panel(slide, left, BODY_TOP, col_w, BODY_H, accent=accent)
            self._panel_label(slide, left + PAD, BODY_TOP + Inches(0.3), col_w - 2 * PAD,
                              label, color=color)
            self._bullets(slide, left + PAD, BODY_TOP + Inches(0.66), col_w - 2 * PAD,
                          BODY_H - Inches(0.95), items or ['Add guidance.'],
                          max_pt=13, min_pt=9, bullet_color=color)
        self._notes(slide, 'These rules keep the conversation about value, not about the rival.')

    def slide_pricing(self):
        pricing = self.card.get('pricing', {})
        notes = pricing.get('notes') or []
        if not any((pricing.get('ia_model'), pricing.get('competitor_model'), notes)):
            return
        slide = self._page('Pricing and packaging')
        col_w = (CONTENT_W - GUTTER) / 2
        panel_h = BODY_H - Inches(1.9)
        pairs = (
            ('Impact Analytics', pricing.get('ia_model', ''), IMPACT_BLUE, WHITE),
            (self.meta['competitor'], pricing.get('competitor_model', ''), WHITE, BLACK),
        )
        for index, (label, text, fill, ink) in enumerate(pairs):
            left = MARGIN + index * (col_w + GUTTER)
            self._panel(slide, left, BODY_TOP, col_w, panel_h, fill=fill,
                        border=None if index == 0 else GRAY_1,
                        accent=None if index == 0 else self.theme.accent)
            self._panel_label(slide, left + PAD, BODY_TOP + Inches(0.26), col_w - 2 * PAD,
                              label, color=mix(WHITE, IMPACT_BLUE, 0.35) if index == 0 else IMPACT_BLUE)
            self._paragraph_block(slide, left + PAD, BODY_TOP + Inches(0.6), col_w - 2 * PAD,
                                  panel_h - Inches(0.85), text or 'Add detail.',
                                  max_pt=13, min_pt=9, color=ink)
        note_top = BODY_TOP + panel_h + Inches(0.2)
        self._panel(slide, MARGIN, note_top, CONTENT_W, BODY_H - panel_h - Inches(0.2),
                    fill=WHITE, accent=GRAY_2)
        self._panel_label(slide, MARGIN + PAD, note_top + Inches(0.2), CONTENT_W - 2 * PAD,
                          'Ground rules', color=GRAY_3)
        self._bullets(slide, MARGIN + PAD, note_top + Inches(0.5), CONTENT_W - 2 * PAD,
                      BODY_H - panel_h - Inches(0.78),
                      notes or ['Record only what the buyer states or the vendor publishes.'],
                      max_pt=11, min_pt=8, bullet_color=GRAY_3)
        self._notes(slide, 'Never quote a rival price you cannot source.')

    def slide_next_steps(self):
        steps = self.card.get('next_steps') or []
        resources = self.card.get('resources') or []
        if not steps and not resources:
            return
        slide = self._page('Next steps and resources')
        left_w = Inches(7.6)
        self._panel(slide, MARGIN, BODY_TOP, left_w, BODY_H, accent=self.theme.accent)
        self._panel_label(slide, MARGIN + PAD, BODY_TOP + Inches(0.3), left_w - 2 * PAD,
                          'Move the deal forward')
        step_top = BODY_TOP + Inches(0.7)
        step_h = min(Inches(1.0), (BODY_H - Inches(1.0)) / max(1, len(steps) or 1))
        for index, step in enumerate(steps[:5]):
            top = step_top + index * step_h
            marker = add_textbox(slide, MARGIN + PAD, top, Inches(0.4), Inches(0.4))
            write_paragraph(marker.text_frame, str(index + 1), Pt(15), bold=True,
                            color=self.theme.accent, font=self.theme.heading_font,
                            space_after=0, line_spacing=1.0, first=True)
            self._paragraph_block(slide, MARGIN + PAD + Inches(0.45), top,
                                  left_w - 2 * PAD - Inches(0.45), step_h - Inches(0.12),
                                  step, max_pt=13, min_pt=9)

        right_left = MARGIN + left_w + GUTTER
        right_w = CONTENT_W - left_w - GUTTER
        self._panel(slide, right_left, BODY_TOP, right_w, BODY_H, fill=WHITE)
        self._panel_label(slide, right_left + PAD, BODY_TOP + Inches(0.3), right_w - 2 * PAD,
                          'Resources')
        top = BODY_TOP + Inches(0.68)
        for row in resources[:8]:
            label = row.get('label') or row.get('url')
            url = row.get('url')
            box = add_textbox(slide, right_left + PAD, top, right_w - 2 * PAD, Inches(0.3))
            p = box.text_frame.paragraphs[0]
            p.line_spacing = 1.1
            run = p.add_run()
            run.text = label
            set_run_font(run, Pt(11), False, IMPACT_BLUE if url else BLACK, self.theme.body_font)
            if url:
                run.hyperlink.address = url
            top += Inches(0.34)
        owner = self.meta.get('owner')
        if owner:
            self._paragraph_block(slide, right_left + PAD, BODY_TOP + BODY_H - Inches(0.6),
                                  right_w - 2 * PAD, Inches(0.4),
                                  'Card owner: %s' % owner, max_pt=10, min_pt=8, color=GRAY_3)
        self._notes(slide, 'Agree the success metric and the readout date in writing.')

    def slide_one_pager(self):
        slide = self._page('One page summary',
                           kicker='Print this  ·  %s' % self.meta['competitor'])
        col_w = (CONTENT_W - GUTTER * 3) / 4
        quadrants = (
            ('Why we win', [row.get('title', '') for row in self.card.get('our_advantages', [])][:4]),
            ('Where they win', (self.card.get('their_strengths') or [])[:4]),
            ('Their gaps', (self.card.get('their_weaknesses') or [])[:4]),
            ('Top objections', self._objection_pairs()),
        )
        panel_h = BODY_H - Inches(1.55)
        for index, (label, items) in enumerate(quadrants):
            left = MARGIN + index * (col_w + GUTTER)
            self._panel(slide, left, BODY_TOP, col_w, panel_h, accent=self.theme.accent)
            self._panel_label(slide, left + PAD, BODY_TOP + Inches(0.24), col_w - 2 * PAD, label)
            self._bullets(slide, left + PAD, BODY_TOP + Inches(0.58), col_w - 2 * PAD,
                          panel_h - Inches(0.82),
                          [item for item in items if item] or ['Add content.'],
                          max_pt=11, min_pt=7.5, bullet_color=IMPACT_BLUE)

        strip_top = BODY_TOP + panel_h + Inches(0.2)
        strip_h = BODY_H - panel_h - Inches(0.2)
        self._panel(slide, MARGIN, strip_top, CONTENT_W, strip_h, fill=IMPACT_BLUE, border=None)
        self._panel_label(slide, MARGIN + PAD, strip_top + Inches(0.18), CONTENT_W - 2 * PAD,
                          'Say this first', color=mix(WHITE, IMPACT_BLUE, 0.35))
        line = (self.card.get('talk_track', {}).get('positioning')
                or self.meta.get('headline')
                or 'Lead with the outcome the buyer needs this season.')
        self._paragraph_block(slide, MARGIN + PAD, strip_top + Inches(0.48),
                              CONTENT_W - 2 * PAD, strip_h - Inches(0.66), line,
                              max_pt=15, min_pt=10, color=WHITE)
        self._notes(slide, 'Print this slide. It is the card a seller carries into the room.')

    def _objection_pairs(self):
        """Compact objection and answer pairs for the printable summary."""
        pairs = []
        for row in (self.card.get('objections') or [])[:2]:
            if row.get('objection'):
                pairs.append('They say: %s' % row['objection'])
            if row.get('response'):
                pairs.append('We say: %s' % row['response'])
        return pairs


def build_presentation(card: dict) -> Presentation:
    return BattlecardDeck(card).build()
