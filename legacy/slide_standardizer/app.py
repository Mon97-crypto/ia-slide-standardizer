import os
import uuid
import io
import zipfile
import copy
import re
from datetime import datetime
from flask import Flask, render_template, request, send_file, jsonify, url_for
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')
app.config['OUTPUT_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'outputs')

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

# ─── Brand Constants ────────────────────────────────────────────────────────────
BRAND = {
    'font': 'Inter Tight',
    'font_fallback': 'Calibri',
    'heading_size': Pt(24),
    'subheading_size': Pt(18),
    'body_size': Pt(14),
    'caption_size': Pt(10),
    'stat_size': Pt(40),
    'slide_width': Emu(12192000),   # 13.33 inches (widescreen)
    'slide_height': Emu(6858000),   # 7.50 inches
    'bg_white': RGBColor(0xFF, 0xFF, 0xFF),
    'bg_dark': RGBColor(0x1C, 0x1B, 0x1B),
    'primary_blue': RGBColor(0x26, 0x4C, 0xD7),
    'accent_blue': RGBColor(0x3A, 0x5E, 0xE3),
    'accent_green': RGBColor(0x3D, 0xD4, 0x99),
    'dark_navy': RGBColor(0x17, 0x37, 0x5E),
    'text_dark': RGBColor(0x1C, 0x1B, 0x1B),
    'text_white': RGBColor(0xFF, 0xFF, 0xFF),
    'text_blue': RGBColor(0x26, 0x4C, 0xD7),
    'light_gray': RGBColor(0xF4, 0xF4, 0xF6),
    'medium_gray': RGBColor(0xD5, 0xDB, 0xE5),
    # Product color mapping (from the image with colored outlines)
    'product_colors': {
        'PlanSmart': RGBColor(0xC4, 0xB5, 0xE3),      # Purple/lavender
        'AssortSmart': RGBColor(0xC4, 0xB5, 0xE3),
        'VisualSmart': RGBColor(0xC4, 0xB5, 0xE3),
        'SizeSmart': RGBColor(0xC4, 0xB5, 0xE3),
        'ItemSmart': RGBColor(0xC4, 0xB5, 0xE3),
        'StoreSmart': RGBColor(0xC4, 0xB5, 0xE3),
        'InventorySmart': RGBColor(0xD4, 0xE1, 0x57),  # Yellow-green
        'SourceSmart': RGBColor(0xD4, 0xE1, 0x57),
        'SpaceSmart': RGBColor(0xD4, 0xE1, 0x57),
        'PriceSmart': RGBColor(0x00, 0xD1, 0xB2),      # Teal/cyan
        'BaseSmart': RGBColor(0x00, 0xD1, 0xB2),
        'PromoSmart': RGBColor(0x00, 0xD1, 0xB2),
        'MarkSmart': RGBColor(0x00, 0xD1, 0xB2),
        'TradeSmart': RGBColor(0x00, 0xD1, 0xB2),
        'MondaySmart': RGBColor(0x26, 0x4C, 0xD7),     # Blue
        'AttributeSmart': RGBColor(0x26, 0x4C, 0xD7),
        'TestSmart': RGBColor(0x26, 0x4C, 0xD7),
    }
}

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static', 'images', 'ia_logo.png')


# ─── Slide Builder Helpers ──────────────────────────────────────────────────────

def set_font(run, size=None, bold=False, color=None, font_name=None):
    """Apply brand font styling to a run."""
    run.font.name = font_name or BRAND['font']
    if size:
        run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_logo(slide, prs):
    """Add IA logo to top-right corner."""
    if os.path.exists(LOGO_PATH):
        logo_w = Emu(1200000)  # ~1.3 inches
        logo_h = Emu(393750)   # proportional
        left = prs.slide_width - logo_w - Emu(200000)
        top = Emu(150000)
        slide.shapes.add_picture(LOGO_PATH, left, top, logo_w, logo_h)


def add_footer_bar(slide, prs, text="", color=None):
    """Add colored footer bar at bottom of slide."""
    bar_h = Emu(300000)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), prs.slide_height - bar_h,
        prs.slide_width, bar_h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color or BRAND['primary_blue']
    bar.line.fill.background()
    if text:
        tf = bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        set_font(run, BRAND['caption_size'], color=BRAND['text_white'])


def make_rounded_rect(slide, left, top, width, height, fill_color, corner_radius=Emu(100000)):
    """Create a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Set corner radius via XML
    try:
        sp = shape._element
        prstGeom = sp.find(qn('a:prstGeom'))
        if prstGeom is None:
            spPr = sp.find(qn('p:spPr'))
            if spPr is not None:
                prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = prstGeom.makeelement(qn('a:avLst'), {})
                prstGeom.append(avLst)
    except:
        pass
    return shape


def add_text_to_shape(shape, text, size=None, bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name=None):
    """Add text to a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size or BRAND['body_size'], bold, color or BRAND['text_dark'], font_name)
    return tf


def add_paragraph(text_frame, text, size=None, bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name=None):
    """Add a new paragraph to existing text frame."""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size or BRAND['body_size'], bold, color or BRAND['text_dark'], font_name)
    return p


# ─── Slide Builders ────────────────────────────────────────────────────────────

def build_intro_slide(prs, title="", subtitle="", date_text="", product_name=""):
    """Build a branded intro/title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Dark background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND['bg_dark']

    # Product accent stripe on left
    product_color = BRAND['product_colors'].get(product_name, BRAND['primary_blue'])
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), Emu(180000), prs.slide_height
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = product_color
    stripe.line.fill.background()

    # Logo
    if os.path.exists(LOGO_PATH):
        logo_w = Emu(2000000)
        logo_h = Emu(656250)
        left = (prs.slide_width - logo_w) // 2
        top = Emu(1200000)
        slide.shapes.add_picture(LOGO_PATH, left, top, logo_w, logo_h)

    # Title
    txBox = slide.shapes.add_textbox(
        Emu(900000), Emu(2400000),
        prs.slide_width - Emu(1800000), Emu(1200000)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title or "Presentation Title"
    set_font(run, Pt(32), bold=True, color=BRAND['text_white'])

    # Subtitle
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        run2 = p2.add_run()
        run2.text = subtitle
        set_font(run2, Pt(18), color=BRAND['text_white'])

    # Date
    if date_text:
        dtBox = slide.shapes.add_textbox(
            Emu(900000), Emu(4200000),
            prs.slide_width - Emu(1800000), Emu(400000)
        )
        p3 = dtBox.text_frame.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = date_text
        set_font(run3, Pt(14), color=RGBColor(0xAA, 0xAA, 0xAA))

    add_footer_bar(slide, prs, color=product_color)
    return slide


def build_section_divider(prs, section_title="", section_number="", product_name=""):
    """Build a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    product_color = BRAND['product_colors'].get(product_name, BRAND['primary_blue'])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = product_color

    # Section number
    if section_number:
        numBox = slide.shapes.add_textbox(
            Emu(600000), Emu(1500000),
            Emu(2000000), Emu(1200000)
        )
        p = numBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = section_number
        set_font(run, Pt(72), bold=True, color=BRAND['text_white'])

    # Section title
    titleBox = slide.shapes.add_textbox(
        Emu(600000), Emu(2800000),
        prs.slide_width - Emu(1200000), Emu(1500000)
    )
    tf = titleBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_title or "Section Title"
    set_font(run, Pt(36), bold=True, color=BRAND['text_white'])

    add_logo(slide, prs)
    return slide


def build_content_slide(prs, title="", body_points=None, product_name=""):
    """Build a standard content slide with title + bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND['bg_white']

    product_color = BRAND['product_colors'].get(product_name, BRAND['primary_blue'])

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        prs.slide_width, Emu(900000)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BRAND['light_gray']
    title_bar.line.fill.background()

    # Title text
    titleBox = slide.shapes.add_textbox(
        Emu(500000), Emu(200000),
        prs.slide_width - Emu(2000000), Emu(600000)
    )
    tf = titleBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title or "Slide Title"
    set_font(run, BRAND['heading_size'], bold=True, color=BRAND['text_dark'])

    # Accent line under title
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(500000), Emu(880000),
        Emu(1500000), Emu(40000)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = product_color
    accent.line.fill.background()

    # Body content
    if body_points:
        bodyBox = slide.shapes.add_textbox(
            Emu(500000), Emu(1100000),
            prs.slide_width - Emu(1000000), prs.slide_height - Emu(1700000)
        )
        tf = bodyBox.text_frame
        tf.word_wrap = True
        for i, point in enumerate(body_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(6)
            p.space_after = Pt(4)
            run = p.add_run()
            run.text = f"• {point}" if not point.startswith('•') else point
            set_font(run, BRAND['body_size'], color=BRAND['text_dark'])

    add_logo(slide, prs)
    add_footer_bar(slide, prs, color=product_color)
    return slide


def build_case_study_slide(prs, title="", challenge="", solution="", impact_stat="",
                           impact_label="", extra_impacts=None, tag="", product_name=""):
    """Build a 3-column case study slide (Challenge | Solution | Impact)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND['bg_white']

    product_color = BRAND['product_colors'].get(product_name, BRAND['primary_blue'])

    # Title bar
    titleBox = slide.shapes.add_textbox(
        Emu(350000), Emu(150000),
        prs.slide_width - Emu(700000), Emu(900000)
    )
    tf = titleBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title or "Case Study Title"
    set_font(run, Pt(20), bold=True, color=BRAND['text_dark'])

    # Tag pill (top right)
    if tag:
        tag_shape = make_rounded_rect(
            slide,
            prs.slide_width - Emu(2500000), Emu(0),
            Emu(2200000), Emu(280000),
            product_color
        )
        add_text_to_shape(tag_shape, tag, Pt(10), bold=True,
                         color=BRAND['text_white'], alignment=PP_ALIGN.CENTER)

    # Three columns
    col_w = Emu(3616500)
    col_h = Emu(4800000)
    col_top = Emu(1300000)
    gap = Emu(200000)
    col_starts = [Emu(440000), Emu(440000) + col_w + gap, Emu(440000) + 2 * (col_w + gap)]

    headers = ['Challenge', 'Solution', 'Impact']
    header_icons = ['⚠', '✦', '▲']
    contents = [challenge, solution, '']

    for i in range(3):
        # Column background
        col_bg = make_rounded_rect(
            slide, col_starts[i], col_top,
            col_w, col_h,
            BRAND['bg_white']
        )
        col_bg.line.color.rgb = BRAND['medium_gray']
        col_bg.line.width = Pt(1)

        # Column header
        header_shape = make_rounded_rect(
            slide,
            col_starts[i] + Emu(70000), col_top + Emu(70000),
            col_w - Emu(140000), Emu(420000),
            BRAND['light_gray']
        )
        add_text_to_shape(header_shape, headers[i], Pt(16), bold=True,
                         color=BRAND['text_dark'], alignment=PP_ALIGN.CENTER)

        # Content
        if i < 2 and contents[i]:
            content_box = slide.shapes.add_textbox(
                col_starts[i] + Emu(120000),
                col_top + Emu(600000),
                col_w - Emu(240000),
                col_h - Emu(700000)
            )
            ctf = content_box.text_frame
            ctf.word_wrap = True
            lines = contents[i].split('\n') if isinstance(contents[i], str) else contents[i] if isinstance(contents[i], list) else [str(contents[i])]
            for j, line in enumerate(lines):
                if j == 0:
                    cp = ctf.paragraphs[0]
                else:
                    cp = ctf.add_paragraph()
                cp.space_before = Pt(6)
                cp.space_after = Pt(4)
                crun = cp.add_run()
                crun.text = f"• {line}" if not line.startswith('•') else line
                set_font(crun, Pt(12), color=BRAND['text_dark'])

    # Impact column special content
    if impact_stat:
        stat_box = slide.shapes.add_textbox(
            col_starts[2] + Emu(120000),
            col_top + Emu(600000),
            col_w - Emu(240000),
            Emu(800000)
        )
        stf = stat_box.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = impact_stat
        set_font(run, BRAND['stat_size'], bold=True, color=product_color)

        if impact_label:
            p2 = stf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = impact_label
            set_font(r2, Pt(12), color=BRAND['text_dark'])

    if extra_impacts:
        y_offset = col_top + Emu(1800000)
        for ei in extra_impacts:
            ei_box = make_rounded_rect(
                slide,
                col_starts[2] + Emu(120000), y_offset,
                col_w - Emu(240000), Emu(500000),
                BRAND['light_gray']
            )
            add_text_to_shape(ei_box, ei, Pt(11), color=BRAND['text_dark'], alignment=PP_ALIGN.CENTER)
            y_offset += Emu(550000)

    add_logo(slide, prs)
    add_footer_bar(slide, prs, color=product_color)
    return slide


def build_content_image_slide(prs, title="", body_text="", image_path=None, product_name=""):
    """Build a two-column slide: text left, image right."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND['bg_white']

    product_color = BRAND['product_colors'].get(product_name, BRAND['primary_blue'])

    # Title
    titleBox = slide.shapes.add_textbox(
        Emu(500000), Emu(200000),
        prs.slide_width - Emu(1000000), Emu(700000)
    )
    p = titleBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title or "Slide Title"
    set_font(run, BRAND['heading_size'], bold=True, color=BRAND['text_dark'])
    titleBox.text_frame.word_wrap = True

    # Left column - text
    half_w = (prs.slide_width - Emu(1200000)) // 2
    textBox = slide.shapes.add_textbox(
        Emu(500000), Emu(1100000),
        half_w, prs.slide_height - Emu(1700000)
    )
    tf = textBox.text_frame
    tf.word_wrap = True
    lines = body_text.split('\n') if body_text else ["Content goes here"]
    for i, line in enumerate(lines):
        if i == 0:
            cp = tf.paragraphs[0]
        else:
            cp = tf.add_paragraph()
        cp.space_before = Pt(4)
        crun = cp.add_run()
        crun.text = line
        set_font(crun, BRAND['body_size'], color=BRAND['text_dark'])

    # Right column - image placeholder or actual image
    img_left = Emu(500000) + half_w + Emu(200000)
    img_w = half_w
    img_h = prs.slide_height - Emu(1700000)
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, img_left, Emu(1100000), img_w, img_h)
    else:
        placeholder = make_rounded_rect(
            slide, img_left, Emu(1100000), img_w, img_h,
            BRAND['light_gray']
        )
        add_text_to_shape(placeholder, "[Image Placeholder]", Pt(16),
                         color=BRAND['medium_gray'], alignment=PP_ALIGN.CENTER)

    add_logo(slide, prs)
    add_footer_bar(slide, prs, color=product_color)
    return slide


# ─── PPTX Processing Engine ────────────────────────────────────────────────────

def extract_slide_content(slide):
    """Extract all text content from a slide."""
    content = {
        'texts': [],
        'images': [],
        'layout_name': slide.slide_layout.name if slide.slide_layout else 'Unknown'
    }

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    is_bold = any(run.font.bold for run in para.runs if run.font.bold)
                    font_size = None
                    for run in para.runs:
                        if run.font.size:
                            font_size = run.font.size
                            break
                    content['texts'].append({
                        'text': text,
                        'bold': is_bold,
                        'size': font_size,
                        'shape_name': shape.name,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    })

        if hasattr(shape, 'image'):
            try:
                img_data = shape.image.blob
                content['images'].append({
                    'data': img_data,
                    'content_type': shape.image.content_type,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })
            except:
                pass

    return content


def classify_slide(content):
    """Classify a slide type based on its content."""
    texts = content['texts']
    if not texts:
        return 'blank'

    text_count = len(texts)
    has_images = len(content['images']) > 0
    all_text = ' '.join(t['text'].lower() for t in texts)

    # Check for case study pattern
    if any(w in all_text for w in ['challenge', 'solution', 'impact']):
        return 'case_study'

    # Check for section divider (very few words, large font)
    if text_count <= 2:
        avg_size = sum(t['size'] or 0 for t in texts) / max(text_count, 1)
        if avg_size > Pt(28):
            return 'section_divider'

    # Title slide detection
    if text_count <= 3 and any(w in all_text for w in ['date', '2025', '2026', 'presented', 'confidential']):
        return 'intro'

    # Check if mostly text (agenda, content)
    if text_count <= 2 and not has_images:
        return 'intro'

    if has_images and text_count > 1:
        return 'content_image'

    return 'content'


def detect_product(content):
    """Detect which IA product is referenced in the slide."""
    all_text = ' '.join(t['text'] for t in content['texts'])
    for product in BRAND['product_colors']:
        if product.lower() in all_text.lower():
            return product
    return ''


def standardize_pptx(input_path, template_type='auto'):
    """Process an input PPTX and output in standardized IA template."""
    src = Presentation(input_path)

    # Create output presentation with standard dimensions
    output = Presentation()
    output.slide_width = BRAND['slide_width']
    output.slide_height = BRAND['slide_height']

    # Ensure we have a blank layout
    blank_layout = output.slide_layouts[6]  # Usually blank

    for slide_idx, slide in enumerate(src.slides):
        content = extract_slide_content(slide)
        slide_type = template_type if template_type != 'auto' else classify_slide(content)
        product = detect_product(content)

        texts = content['texts']
        title_text = texts[0]['text'] if texts else f"Slide {slide_idx + 1}"
        body_texts = [t['text'] for t in texts[1:]] if len(texts) > 1 else []

        if slide_type == 'intro':
            subtitle = body_texts[0] if body_texts else ""
            date = body_texts[1] if len(body_texts) > 1 else ""
            build_intro_slide(output, title_text, subtitle, date, product)

        elif slide_type == 'section_divider':
            num = str(slide_idx + 1).zfill(2)
            build_section_divider(output, title_text, num, product)

        elif slide_type == 'case_study':
            # Parse challenge/solution/impact from body
            challenge_lines = []
            solution_lines = []
            impact_lines = []
            current_section = None
            for t in body_texts:
                lower = t.lower()
                if 'challenge' in lower:
                    current_section = 'challenge'
                    continue
                elif 'solution' in lower:
                    current_section = 'solution'
                    continue
                elif 'impact' in lower:
                    current_section = 'impact'
                    continue
                if current_section == 'challenge':
                    challenge_lines.append(t)
                elif current_section == 'solution':
                    solution_lines.append(t)
                elif current_section == 'impact':
                    impact_lines.append(t)
                else:
                    challenge_lines.append(t)

            stat = impact_lines[0] if impact_lines else ""
            label = impact_lines[1] if len(impact_lines) > 1 else ""
            extras = impact_lines[2:] if len(impact_lines) > 2 else []

            build_case_study_slide(
                output, title_text,
                '\n'.join(challenge_lines),
                '\n'.join(solution_lines),
                stat, label, extras,
                product_name=product
            )

        elif slide_type == 'content_image':
            body = '\n'.join(body_texts)
            # Save first image temporarily
            img_path = None
            if content['images']:
                img_data = content['images'][0]
                ext = img_data['content_type'].split('/')[-1]
                if ext == 'jpeg':
                    ext = 'jpg'
                img_path = os.path.join(app.config['UPLOAD_FOLDER'], f"temp_img_{slide_idx}.{ext}")
                with open(img_path, 'wb') as f:
                    f.write(img_data['data'])
            build_content_image_slide(output, title_text, body, img_path, product)

        else:  # content
            build_content_slide(output, title_text, body_texts if body_texts else [title_text], product)

    return output


def process_screenshot(image_path, slide_type='content'):
    """Process a screenshot image into a PPTX slide using OCR."""
    output = Presentation()
    output.slide_width = BRAND['slide_width']
    output.slide_height = BRAND['slide_height']

    # Try OCR
    extracted_text = ""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(image_path)
        extracted_text = pytesseract.image_to_string(img)
    except Exception as e:
        extracted_text = f"[OCR unavailable - image placed directly]\nError: {str(e)}"

    lines = [l.strip() for l in extracted_text.split('\n') if l.strip()]
    title = lines[0] if lines else "Extracted Content"
    body = lines[1:] if len(lines) > 1 else ["Content extracted from screenshot"]

    if slide_type == 'intro':
        build_intro_slide(output, title, body[0] if body else "")
    elif slide_type == 'section_divider':
        build_section_divider(output, title)
    elif slide_type == 'case_study':
        build_case_study_slide(output, title, '\n'.join(body[:3]), '\n'.join(body[3:6]),
                              body[6] if len(body) > 6 else "")
    elif slide_type == 'content_image':
        build_content_image_slide(output, title, '\n'.join(body), image_path)
    else:
        build_content_slide(output, title, body)

    return output


# ─── Flask Routes ───────────────────────────────────────────────────────────────

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/process', methods=['POST'])
def process():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    template_type = request.form.get('template_type', 'auto')
    product_name = request.form.get('product_name', '')

    # Save uploaded file
    job_id = str(uuid.uuid4())[:8]
    filename = secure_filename(file.filename)
    upload_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{job_id}_{filename}")
    file.save(upload_path)

    try:
        ext = filename.rsplit('.', 1)[-1].lower()

        if ext in ('pptx',):
            result = standardize_pptx(upload_path, template_type)
        elif ext in ('png', 'jpg', 'jpeg', 'bmp', 'tiff', 'webp'):
            result = process_screenshot(upload_path, template_type)
        else:
            return jsonify({'error': f'Unsupported file type: .{ext}'}), 400

        # Save output
        output_name = f"IA_Standardized_{job_id}.pptx"
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_name)
        result.save(output_path)

        return jsonify({
            'success': True,
            'filename': output_name,
            'download_url': f'/download/{output_name}',
            'slide_count': len(result.slides)
        })

    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'traceback': traceback.format_exc()}), 500
    finally:
        # Cleanup upload
        try:
            os.remove(upload_path)
        except:
            pass


@app.route('/download/<filename>')
def download(filename):
    filepath = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True, download_name=filename)
    return jsonify({'error': 'File not found'}), 404


@app.route('/build-custom', methods=['POST'])
def build_custom():
    """Build a custom slide deck from JSON specification."""
    data = request.get_json()
    if not data:
        return jsonify({'error': 'No data provided'}), 400

    slides_spec = data.get('slides', [])
    product_name = data.get('product_name', '')

    output = Presentation()
    output.slide_width = BRAND['slide_width']
    output.slide_height = BRAND['slide_height']

    for spec in slides_spec:
        stype = spec.get('type', 'content')
        if stype == 'intro':
            build_intro_slide(output, spec.get('title', ''), spec.get('subtitle', ''),
                            spec.get('date', ''), product_name)
        elif stype == 'section_divider':
            build_section_divider(output, spec.get('title', ''), spec.get('number', ''), product_name)
        elif stype == 'content':
            build_content_slide(output, spec.get('title', ''),
                              spec.get('points', []), product_name)
        elif stype == 'case_study':
            build_case_study_slide(output, spec.get('title', ''),
                                  spec.get('challenge', ''), spec.get('solution', ''),
                                  spec.get('impact_stat', ''), spec.get('impact_label', ''),
                                  spec.get('extra_impacts', []),
                                  spec.get('tag', ''), product_name)
        elif stype == 'content_image':
            build_content_image_slide(output, spec.get('title', ''),
                                    spec.get('body', ''), None, product_name)

    job_id = str(uuid.uuid4())[:8]
    output_name = f"IA_Custom_{job_id}.pptx"
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_name)
    output.save(output_path)

    return jsonify({
        'success': True,
        'filename': output_name,
        'download_url': f'/download/{output_name}',
        'slide_count': len(output.slides)
    })


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=True)
