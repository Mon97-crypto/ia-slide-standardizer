"""Render a sales deck as an Impact Analytics branded .pptx.

The model decides what the deck says; this module decides how it looks. Layout
is computed here rather than described by the model, because a model asked to
place text boxes produces slides that overlap on the first long sentence.

Colours and type follow the Impact Analytics brand: Impact Blue for emphasis,
off-white grounds, Inter Tight throughout, with Spectral standing in for the
licensed display serif when it is unavailable on the rendering machine.
"""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

BLUE = RGBColor(0x26, 0x4C, 0xD7)
BLUE_DEEP = RGBColor(0x16, 0x26, 0x7A)
OFF_WHITE = RGBColor(0xF4, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1C, 0x1B, 0x1B)
GREY = RGBColor(0x8A, 0x88, 0x84)
LINE = RGBColor(0xE4, 0xE4, 0xE4)
ORANGE = RGBColor(0xFF, 0x6F, 0x1C)

BODY_FONT = "Inter Tight"
DISPLAY_FONT = "Spectral"

# 16:9 at the size PowerPoint uses for widescreen.
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN = Emu(838200)          # 0.9 inch
CONTENT_W = SLIDE_W - MARGIN * 2

# The deck the model fills in. Types are fixed so layout stays predictable.
DECK_SCHEMA = {
    "type": "object",
    "properties": {
        "title": {"type": "string", "description": "Deck title, under 60 characters."},
        "subtitle": {"type": "string", "description": "One line naming the audience and the ask."},
        "slides": {
            "type": "array",
            "minItems": 1,
            "items": {
                "type": "object",
                "properties": {
                    "type": {"type": "string",
                             "enum": ["section", "bullets", "comparison", "stats", "quote", "close"]},
                    "heading": {"type": "string"},
                    "subheading": {"type": "string"},
                    "bullets": {"type": "array", "items": {"type": "string"},
                                "description": "Used by bullets and close. Five at most, one line each."},
                    "left_title": {"type": "string", "description": "comparison only."},
                    "right_title": {"type": "string", "description": "comparison only."},
                    "left": {"type": "array", "items": {"type": "string"}},
                    "right": {"type": "array", "items": {"type": "string"}},
                    "stats": {
                        "type": "array",
                        "description": "stats only. Three at most.",
                        "items": {
                            "type": "object",
                            "properties": {
                                "value": {"type": "string", "description": "Short, such as 9 to 18 months."},
                                "label": {"type": "string"},
                            },
                            "required": ["value", "label"],
                            "additionalProperties": False,
                        },
                    },
                    "quote": {"type": "string", "description": "quote only."},
                    "attribution": {"type": "string"},
                    "note": {"type": "string", "description": "Speaker note for this slide."},
                },
                "required": ["type", "heading"],
                "additionalProperties": False,
            },
        },
    },
    "required": ["title", "subtitle", "slides"],
    "additionalProperties": False,
}


def _text(frame, text: str, size: int, *, bold: bool = False,
          color: RGBColor = BLACK, font: str = BODY_FONT,
          align=PP_ALIGN.LEFT, space_after: int = 0) -> None:
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.alignment = align
    para.space_after = Pt(space_after)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height).text_frame


def _blank(deck: Presentation, background: RGBColor = WHITE):
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = background
    return slide


def _rule(slide, top, width=Emu(600000), color: RGBColor = BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top, width, Emu(48000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


def _bullet_list(slide, items: list[str], top, width, size: int = 15,
                 color: RGBColor = BLACK) -> None:
    """Bullets as separate boxes, so one long line cannot push the rest off."""
    y = top
    for item in items[:6]:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, MARGIN + Emu(20000), y + Emu(70000),
            Emu(90000), Emu(90000))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        dot.shadow.inherit = False
        frame = _box(slide, MARGIN + Emu(220000), y, width - Emu(220000), Emu(520000))
        _text(frame, item, size, color=color)
        y += Emu(620000)


def _speaker_note(slide, note: str) -> None:
    if note:
        slide.notes_slide.notes_text_frame.text = note


def _title_slide(deck: Presentation, spec: dict[str, Any]) -> None:
    slide = _blank(deck, BLUE)
    # A soft block behind the wordmark keeps the type off a flat field.
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), SLIDE_H - Emu(1500000),
                                   SLIDE_W, Emu(1500000))
    panel.fill.solid()
    panel.fill.fore_color.rgb = BLUE_DEEP
    panel.line.fill.background()
    panel.shadow.inherit = False

    _text(_box(slide, MARGIN, Emu(1900000), CONTENT_W, Emu(1700000)),
          spec.get("title", "Impact Analytics"), 40, color=WHITE, font=DISPLAY_FONT)
    _text(_box(slide, MARGIN, Emu(3500000), CONTENT_W, Emu(900000)),
          spec.get("subtitle", ""), 17, color=OFF_WHITE)
    _text(_box(slide, MARGIN, SLIDE_H - Emu(1100000), CONTENT_W, Emu(500000)),
          "Impact Analytics", 13, bold=True, color=WHITE)


def _section_slide(deck, s):
    slide = _blank(deck, OFF_WHITE)
    _rule(slide, Emu(2600000))
    _text(_box(slide, MARGIN, Emu(2800000), CONTENT_W, Emu(1300000)),
          s.get("heading", ""), 34, font=DISPLAY_FONT)
    if s.get("subheading"):
        _text(_box(slide, MARGIN, Emu(4000000), CONTENT_W, Emu(800000)),
              s["subheading"], 15, color=GREY)
    return slide


def _heading(slide, s, colour: RGBColor = BLACK):
    _text(_box(slide, MARGIN, Emu(700000), CONTENT_W, Emu(900000)),
          s.get("heading", ""), 26, font=DISPLAY_FONT, color=colour)
    if s.get("subheading"):
        _text(_box(slide, MARGIN, Emu(1450000), CONTENT_W, Emu(600000)),
              s["subheading"], 14, color=GREY)
    _rule(slide, Emu(2050000), width=Emu(420000))


def _bullets_slide(deck, s):
    slide = _blank(deck)
    _heading(slide, s)
    _bullet_list(slide, s.get("bullets") or [], Emu(2400000), CONTENT_W)
    return slide


def _comparison_slide(deck, s):
    slide = _blank(deck)
    _heading(slide, s)
    col_w = (CONTENT_W - Emu(400000)) // 2
    for index, (title, items, accent) in enumerate((
            (s.get("left_title", "Impact Analytics"), s.get("left") or [], BLUE),
            (s.get("right_title", "Them"), s.get("right") or [], ORANGE))):
        x = MARGIN + index * (col_w + Emu(400000))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(2350000),
                                      col_w, Emu(3300000))
        card.fill.solid()
        card.fill.fore_color.rgb = OFF_WHITE if index else WHITE
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        card.shadow.inherit = False
        card.text_frame.text = ""

        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Emu(2350000),
                                      col_w, Emu(60000))
        head.fill.solid()
        head.fill.fore_color.rgb = accent
        head.line.fill.background()
        head.shadow.inherit = False

        _text(_box(slide, x + Emu(280000), Emu(2560000), col_w - Emu(560000), Emu(450000)),
              title, 15, bold=True, color=accent)
        y = Emu(3120000)
        for item in items[:5]:
            _text(_box(slide, x + Emu(280000), y, col_w - Emu(560000), Emu(520000)),
                  item, 12, color=BLACK)
            y += Emu(500000)
    return slide


def _stats_slide(deck, s):
    slide = _blank(deck, OFF_WHITE)
    _heading(slide, s)
    stats = (s.get("stats") or [])[:3]
    if not stats:
        return slide
    gap = Emu(300000)
    card_w = (CONTENT_W - gap * (len(stats) - 1)) // len(stats)
    for index, stat in enumerate(stats):
        x = MARGIN + index * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(2500000),
                                      card_w, Emu(2100000))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        card.shadow.inherit = False
        card.text_frame.text = ""
        _text(_box(slide, x + Emu(200000), Emu(2800000), card_w - Emu(400000), Emu(800000)),
              stat.get("value", ""), 34, bold=True, color=BLUE,
              font=DISPLAY_FONT, align=PP_ALIGN.CENTER)
        _text(_box(slide, x + Emu(200000), Emu(3700000), card_w - Emu(400000), Emu(700000)),
              stat.get("label", ""), 12, color=GREY, align=PP_ALIGN.CENTER)
    return slide


def _quote_slide(deck, s):
    slide = _blank(deck, BLUE)
    _text(_box(slide, MARGIN, Emu(2100000), CONTENT_W, Emu(2200000)),
          s.get("quote", s.get("heading", "")), 26, color=WHITE, font=DISPLAY_FONT)
    if s.get("attribution"):
        _text(_box(slide, MARGIN, Emu(4300000), CONTENT_W, Emu(600000)),
              s["attribution"], 13, color=OFF_WHITE)
    return slide


def _close_slide(deck, s):
    slide = _blank(deck)
    _heading(slide, s, colour=BLUE)
    _bullet_list(slide, s.get("bullets") or [], Emu(2400000), CONTENT_W, size=16)
    return slide


BUILDERS = {
    "section": _section_slide,
    "bullets": _bullets_slide,
    "comparison": _comparison_slide,
    "stats": _stats_slide,
    "quote": _quote_slide,
    "close": _close_slide,
}


def build(spec: dict[str, Any]) -> io.BytesIO:
    """Render a deck spec into a .pptx in memory."""
    deck = Presentation()
    deck.slide_width, deck.slide_height = SLIDE_W, SLIDE_H

    _title_slide(deck, spec)
    for entry in spec.get("slides") or []:
        builder = BUILDERS.get(entry.get("type"), _bullets_slide)
        slide = builder(deck, entry)
        if slide is not None:
            _speaker_note(slide, entry.get("note", ""))

    buffer = io.BytesIO()
    deck.save(buffer)
    buffer.seek(0)
    return buffer
