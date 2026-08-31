"""Turn an uploaded file into indexable text, then into retrievable chunks.

Battlecards arrive as PDF and DOCX, client lists as XLSX and CSV, and decks as
PPTX. Accepting only .txt, as the prototype did, rejects almost every real
competitive artefact, so each of those formats is parsed here.

Every extractor is optional at import time. A missing library degrades that one
format to a clear message instead of taking the whole service down.
"""
from __future__ import annotations

import csv
import io
import re
from typing import Any

TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".rst", ".log", ".json"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | {
    ".csv", ".tsv", ".pdf", ".docx", ".pptx", ".xlsx", ".xlsm",
}

# Roughly 1200 characters with 150 of overlap keeps a passage self-contained
# while staying small enough that several fit in one prompt.
CHUNK_CHARS = 1200
CHUNK_OVERLAP = 150


class ExtractionError(Exception):
    """Raised when a file cannot be turned into text."""


def extension_of(filename: str) -> str:
    name = (filename or "").lower().strip()
    dot = name.rfind(".")
    return name[dot:] if dot != -1 else ""


def _decode(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _from_delimited(data: bytes, delimiter: str) -> str:
    text = _decode(data)
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    lines = []
    for row in reader:
        cells = [c.strip() for c in row if c and c.strip()]
        if cells:
            lines.append(" | ".join(cells))
    return "\n".join(lines)


def _from_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - depends on install
        raise ExtractionError("PDF support needs the pypdf package.") from exc
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for number, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            pages.append(f"[page {number}]\n{text.strip()}")
    if not pages:
        raise ExtractionError(
            "No selectable text in this PDF. It is most likely a scan, so add "
            "the key points as a note instead.")
    return "\n\n".join(pages)


def _from_docx(data: bytes) -> str:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("DOCX support needs the python-docx package.") from exc
    document = docx.Document(io.BytesIO(data))
    blocks = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                blocks.append(" | ".join(cells))
    return "\n".join(blocks)


def _from_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("PPTX support needs the python-pptx package.") from exc
    deck = Presentation(io.BytesIO(data))
    slides = []
    for number, slide in enumerate(deck.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if parts:
            slides.append(f"[slide {number}]\n" + "\n".join(parts))
    return "\n\n".join(slides)


def _from_xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("XLSX support needs the openpyxl package.") from exc
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets = []
    for worksheet in workbook.worksheets:
        rows = []
        for row in worksheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c not in (None, "")]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[sheet {worksheet.title}]\n" + "\n".join(rows))
    workbook.close()
    return "\n\n".join(sheets)


_EXTRACTORS = {
    ".csv": lambda d: _from_delimited(d, ","),
    ".tsv": lambda d: _from_delimited(d, "\t"),
    ".pdf": _from_pdf,
    ".docx": _from_docx,
    ".pptx": _from_pptx,
    ".xlsx": _from_xlsx,
    ".xlsm": _from_xlsx,
}


def extract_text(data: bytes, filename: str) -> str:
    """Extract plain text from an uploaded file."""
    if not data:
        raise ExtractionError("That file is empty.")
    ext = extension_of(filename)
    if ext in TEXT_EXTENSIONS:
        return normalise_whitespace(_decode(data))
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ExtractionError(
            f"Cannot read '{ext or 'that file type'}'. Supported: {supported}.")
    try:
        text = extractor(data)
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Could not read this file: {exc}") from exc
    text = normalise_whitespace(text)
    if not text.strip():
        raise ExtractionError("No readable text was found in that file.")
    return text


def normalise_whitespace(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(text: str, size: int = CHUNK_CHARS,
               overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping passages, preferring paragraph boundaries."""
    text = normalise_whitespace(text)
    if not text:
        return []
    if len(text) <= size:
        return [text]

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + size, len(text))
        if end < len(text):
            # Prefer to break on a paragraph, then a sentence, then a space.
            window = text[start:end]
            for marker in ("\n\n", ". ", "\n", " "):
                cut = window.rfind(marker)
                if cut > size * 0.5:
                    end = start + cut + len(marker)
                    break
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        if end >= len(text):
            break
        start = max(end - overlap, start + 1)
    return chunks


def summarise_source(filename: str, text: str) -> dict[str, Any]:
    return {
        "file_name": filename,
        "characters": len(text),
        "words": len(text.split()),
        "chunks": len(chunk_text(text)),
    }
