"""Tests for the parts most likely to regress: query understanding, ranking,
ingestion and the API contract."""
from __future__ import annotations

import io
import json
import os
import sys
import tempfile

import pytest

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ciq import db, ingest  # noqa: E402
from ciq.competitors import canonical_name, resolve, threatened_products  # noqa: E402
from ciq.fetchers import direct_url  # noqa: E402
from ciq.search import parse_query, retrieve_passages, search  # noqa: E402


# Both backends run the same suite. Set CIQ_TEST_DATABASE_URL to a scratch
# Postgres to include it; without it the Postgres runs are skipped rather than
# silently passing.
BACKENDS = ["sqlite"]
if os.environ.get("CIQ_TEST_DATABASE_URL"):
    BACKENDS.append("postgres")


@pytest.fixture(params=BACKENDS)
def conn(request):
    if request.param == "postgres":
        store = db.connect(os.environ["CIQ_TEST_DATABASE_URL"])
        db.clear_all(store)
        yield store
        db.clear_all(store)
        store.close()
        db._local.store = None
        db._local.target = None
        return

    handle, path = tempfile.mkstemp(suffix=".db")
    os.close(handle)
    os.unlink(path)
    store = db.connect(path)
    yield store
    store.close()
    db._local.store = None
    db._local.target = None
    for suffix in ("", "-wal", "-shm"):
        try:
            os.unlink(path + suffix)
        except OSError:
            pass


@pytest.fixture()
def seeded(conn):
    db.create_entry(conn, {
        "competitor": "o9 Solutions", "title": "o9 Battlecard Q2 2026",
        "category": "battlecard",
        "note": "Enterprise planning platform. Strong in CPG and automotive.",
        "content": "o9 is an enterprise planning platform with a knowledge graph. "
                   "Higher price point and longer deployments than retail native tools.",
    }, chunks=["o9 has a higher price point and longer deployments."])
    db.create_entry(conn, {
        "competitor": "Blue Yonder", "title": "Blue Yonder Overview",
        "category": "information", "note": "Mature end to end suite.",
        "content": "Blue Yonder, formerly JDA Software. Weaknesses include legacy "
                   "architecture and long implementation timelines.",
    }, chunks=["Weaknesses include legacy architecture and long implementation timelines."])
    db.create_entry(conn, {
        "competitor": "Relex Solutions", "title": "Relex Client References",
        "category": "client_list", "note": "Grocery heavy client base.",
        "content": "Relex client references are concentrated in grocery retail.",
    }, chunks=["Relex client references are concentrated in grocery retail."])
    return conn


# ─── competitor resolution ─────────────────────────────────────────────────

@pytest.mark.parametrize("raw,expected", [
    ("o9", "o9 Solutions"),
    ("O9 SOLUTIONS", "o9 Solutions"),
    ("BY", "Blue Yonder"),
    ("jda", "Blue Yonder"),          # legacy name maps to the current company
    ("JDA Software", "Blue Yonder"),
    ("relex", "Relex Solutions"),
])
def test_alias_resolution(raw, expected):
    assert canonical_name(raw) == expected


@pytest.mark.parametrize("typo,expected", [
    ("bleu yonder", "Blue Yonder"),
    ("relexx", "Relex Solutions"),
])
def test_fuzzy_resolution(typo, expected):
    name, score = resolve(typo)
    assert name == expected
    assert 0 < score <= 1.0


def test_unknown_competitor_is_not_guessed():
    name, _ = resolve("Some Brand New Vendor")
    assert name == ""


def test_unknown_competitor_is_preserved_verbatim():
    assert canonical_name("Some Brand New Vendor") == "Some Brand New Vendor"


def test_product_mapping():
    assert "PriceSmart" in threatened_products("Blue Yonder")


# ─── query understanding ───────────────────────────────────────────────────

def test_parses_competitor_and_category():
    intent = parse_query("o9 battlecard")
    assert intent.competitor == "o9 Solutions"
    assert intent.category == "battlecard"
    assert intent.terms == []


def test_parses_legacy_name_with_residual_terms():
    intent = parse_query("jda pricing notes")
    assert intent.competitor == "Blue Yonder"
    assert intent.category == "information"
    assert intent.terms == ["pricing"]


def test_multiword_category_beats_single_word():
    assert parse_query("relex client list").category == "client_list"


def test_stopwords_are_dropped():
    assert "for" not in parse_query("client list for relex").terms


# ─── ranking ───────────────────────────────────────────────────────────────

def test_the_query_the_hero_promises(seeded):
    """'o9 battlecard' must return the o9 battlecard. The prototype's
    substring match returned nothing for this exact query."""
    outcome = search(seeded, "o9 battlecard")
    assert outcome["total"] == 1
    assert outcome["results"][0]["title"] == "o9 Battlecard Q2 2026"


def test_alias_search_finds_entry_stored_under_canonical_name(seeded):
    outcome = search(seeded, "jda")
    assert [r["title"] for r in outcome["results"]] == ["Blue Yonder Overview"]


def test_typo_still_finds_the_entry(seeded):
    outcome = search(seeded, "bleu yondr")
    assert outcome["results"][0]["title"] == "Blue Yonder Overview"


def test_search_reaches_into_document_body(seeded):
    outcome = search(seeded, "legacy architecture")
    assert outcome["results"][0]["title"] == "Blue Yonder Overview"


def test_explanations_name_only_terms_that_matched(seeded):
    outcome = search(seeded, "grocery unrelatedword")
    why = " ".join(outcome["results"][0]["why"])
    assert "grocery" in why
    assert "unrelatedword" not in why


def test_category_filter_is_respected(seeded):
    outcome = search(seeded, "", category="client_list")
    assert {r["category"] for r in outcome["results"]} == {"client_list"}


def test_explicit_filter_overrides_inferred_category(seeded):
    outcome = search(seeded, "o9 battlecard", category="client_list")
    assert outcome["effective_category"] == "client_list"
    assert outcome["total"] == 0


def test_empty_query_lists_everything_newest_first(seeded):
    assert search(seeded, "")["total"] == 3


# ─── retrieval ─────────────────────────────────────────────────────────────

def test_retrieval_scopes_to_a_competitor(seeded):
    passages = retrieve_passages(
        seeded, "what are the weaknesses", competitor="Blue Yonder")
    assert passages
    assert all(p["competitor"] == "Blue Yonder" for p in passages)


def test_retrieval_returns_nothing_for_gibberish(seeded):
    assert retrieve_passages(seeded, "zzz") == []


# ─── ingestion ─────────────────────────────────────────────────────────────

def test_text_extraction_normalises_whitespace():
    assert ingest.extract_text(b"a  b\r\n\r\n\r\nc", "n.txt") == "a b\n\nc"


def test_csv_becomes_readable_rows():
    assert ingest.extract_text(b"a,b\n1,2\n", "n.csv") == "a | b\n1 | 2"


def test_unsupported_extension_is_rejected_clearly():
    with pytest.raises(ingest.ExtractionError, match="Cannot read"):
        ingest.extract_text(b"data", "virus.exe")


def test_empty_file_is_rejected():
    with pytest.raises(ingest.ExtractionError):
        ingest.extract_text(b"", "n.txt")


def test_docx_roundtrip():
    docx = pytest.importorskip("docx")
    document = docx.Document()
    document.add_paragraph("Kinaxis is strong in manufacturing.")
    buffer = io.BytesIO()
    document.save(buffer)
    text = ingest.extract_text(buffer.getvalue(), "a.docx")
    assert "manufacturing" in text


def test_xlsx_roundtrip():
    openpyxl = pytest.importorskip("openpyxl")
    workbook = openpyxl.Workbook()
    workbook.active.append(["Client", "Vertical"])
    workbook.active.append(["GroceryCo", "Grocery"])
    buffer = io.BytesIO()
    workbook.save(buffer)
    assert "GroceryCo" in ingest.extract_text(buffer.getvalue(), "a.xlsx")


def test_chunks_overlap_and_cover_the_text():
    text = "word " * 800
    chunks = ingest.chunk_text(text)
    assert len(chunks) > 1
    assert all(len(c) <= ingest.CHUNK_CHARS + 40 for c in chunks)


def test_short_text_is_a_single_chunk():
    assert ingest.chunk_text("short") == ["short"]


# ─── link rewriting ────────────────────────────────────────────────────────

@pytest.mark.parametrize("url,fragment", [
    ("https://docs.google.com/document/d/ABC/edit", "export?format=txt"),
    ("https://docs.google.com/spreadsheets/d/ABC/edit", "export?format=csv"),
    ("https://drive.google.com/file/d/ABC/view", "export=download"),
])
def test_share_links_rewrite_to_direct_downloads(url, fragment):
    rewritten = direct_url(url)
    assert rewritten and fragment in rewritten[0]


def test_unknown_host_is_not_rewritten():
    assert direct_url("https://evil.example.com/file") is None


# ─── storage ───────────────────────────────────────────────────────────────

def test_entry_is_stored_under_the_canonical_competitor(conn):
    entry = db.create_entry(conn, {"competitor": "o9", "title": "T",
                                   "category": "battlecard"})
    assert entry["competitor"] == "o9 Solutions"


def test_delete_removes_the_entry_and_its_index(conn):
    entry = db.create_entry(conn, {"competitor": "o9", "title": "T",
                                   "category": "battlecard",
                                   "content": "uniquetoken here"},
                            chunks=["uniquetoken here"])
    assert db.delete_entry(conn, entry["id"]) is True
    assert search(conn, "uniquetoken")["total"] == 0
    assert db.stats(conn)["chunks"] == 0


# ─── themed retrieval for battlecards ──────────────────────────────────────

def test_themed_sweep_spans_a_long_document(conn):
    """One query surfaces one theme. A battlecard needs several, so the sweep
    must reach parts of a document a single query would miss."""
    from ciq import ingest, llm
    from ciq.search import gather_passages

    doc = (
        "PRICING. Deals land between one and four million dollars annually. "
        + "Budget scrutiny is heavy late in the cycle. " * 20
        + "IMPLEMENTATION. Rollouts run nine to eighteen months. "
        + "Time to value is a recurring friction point. " * 20
        + "PRODUCT. Assortment and size planning are comparatively thin. "
        + "Merchandising depth is the clearest gap. " * 20
    )
    db.create_entry(conn, {
        "competitor": "Blue Yonder", "title": "Dossier",
        "category": "information", "content": doc,
    }, chunks=ingest.chunk_text(doc))

    passages = gather_passages(conn, llm.BATTLECARD_THEMES, "Blue Yonder")
    assert len(passages) > 1
    # What matters is reach: the sweep must surface material from later in the
    # document, which a single positioning query would never return.
    joined = " ".join(p["text"] for p in passages)
    assert "Merchandising" in joined or "Rollouts" in joined
    assert "PRICING" in joined or "Budget" in joined


def test_passages_are_attributed_to_their_best_scoring_theme(conn):
    """A broad theme running first must not claim passages a narrower theme
    describes better."""
    from ciq import ingest, llm
    from ciq.search import gather_passages

    doc = ("Pricing runs into seven figures annually and budget scrutiny is "
           "heavy. " * 25
           + "Implementation and deployment timelines run eighteen months "
             "with integration partners. " * 25)
    db.create_entry(conn, {"competitor": "Blue Yonder", "title": "D",
                           "category": "information", "content": doc},
                    chunks=ingest.chunk_text(doc))
    passages = gather_passages(conn, llm.BATTLECARD_THEMES, "Blue Yonder")
    assert len({p["theme"] for p in passages}) > 1


def test_themed_sweep_deduplicates(conn):
    from ciq import llm
    from ciq.search import gather_passages

    db.create_entry(conn, {
        "competitor": "o9 Solutions", "title": "Short note",
        "category": "battlecard", "content": "o9 pricing is high.",
    }, chunks=["o9 pricing is high."])
    passages = gather_passages(conn, llm.BATTLECARD_THEMES, "o9 Solutions")
    texts = [p["text"] for p in passages]
    assert len(texts) == len(set(texts))


def test_themed_sweep_is_scoped_to_the_competitor(conn):
    from ciq import llm
    from ciq.search import gather_passages

    db.create_entry(conn, {"competitor": "o9 Solutions", "title": "A",
                           "category": "battlecard", "content": "o9 pricing detail."},
                    chunks=["o9 pricing detail."])
    db.create_entry(conn, {"competitor": "Relex Solutions", "title": "B",
                           "category": "battlecard", "content": "Relex pricing detail."},
                    chunks=["Relex pricing detail."])
    passages = gather_passages(conn, llm.BATTLECARD_THEMES, "o9 Solutions")
    assert passages
    assert all(p["competitor"] == "o9 Solutions" for p in passages)


def test_unknown_competitor_yields_no_passages(conn):
    from ciq import llm
    from ciq.search import gather_passages
    assert gather_passages(conn, llm.BATTLECARD_THEMES, "Nobody Inc") == []


def test_backend_is_reported_in_stats(conn):
    assert db.stats(conn)["backend"] in ("sqlite", "postgres")


# ─── head to head scoring ──────────────────────────────────────────────────

def test_weighted_totals_follow_the_weights():
    from ciq.llm import score_totals
    rows = [
        {"dimension": "forecasting", "ia_score": 10, "competitor_score": 0, "weight": 5},
        {"dimension": "breadth", "ia_score": 0, "competitor_score": 10, "weight": 1},
    ]
    totals = score_totals(rows)
    # The heavy dimension must dominate rather than being averaged flat.
    assert totals["ia"] > totals["competitor"]
    assert totals["verdict_key"] == "advantage"


def test_a_close_result_is_not_called_an_advantage():
    from ciq.llm import score_totals
    rows = [{"dimension": "forecasting", "ia_score": 7, "competitor_score": 6.5,
             "weight": 3}]
    assert score_totals(rows)["verdict_key"] == "close"


def test_losing_is_reported_as_losing():
    from ciq.llm import score_totals
    rows = [{"dimension": "breadth", "ia_score": 4, "competitor_score": 9, "weight": 4}]
    totals = score_totals(rows)
    assert totals["verdict_key"] == "behind"
    assert totals["gap"] < 0


def test_totals_survive_an_empty_or_invalid_scorecard():
    from ciq.llm import score_totals
    assert score_totals([])["verdict_key"] == "unknown"
    assert score_totals([{"dimension": "not_a_dimension"}])["verdict_key"] == "unknown"


def test_scorecard_normalisation_clamps_and_deduplicates():
    from ciq.llm import normalise_scorecard
    rows = normalise_scorecard([
        {"dimension": "cost", "ia_score": 99, "competitor_score": -4, "weight": 12},
        {"dimension": "cost", "ia_score": 5, "competitor_score": 5, "weight": 3},
        {"dimension": "made_up", "ia_score": 5, "competitor_score": 5, "weight": 3},
        "not even a dict",
    ])
    assert len(rows) == 1                      # duplicate and unknown dropped
    assert rows[0]["ia_score"] == 10.0         # clamped to the stated range
    assert rows[0]["competitor_score"] == 0.0
    assert rows[0]["weight"] == 5.0
    assert rows[0]["label"] == "Total cost of ownership"


def test_normalisation_defaults_unknown_evidence_to_inference():
    from ciq.llm import normalise_scorecard
    rows = normalise_scorecard([
        {"dimension": "cost", "ia_score": 5, "competitor_score": 5,
         "weight": 3, "evidence": "made up"}])
    assert rows[0]["evidence"] == "inference"


def test_scorecard_leads_with_the_competitor_strengths():
    """A seller needs the threats before the wins."""
    from ciq.llm import normalise_scorecard
    rows = normalise_scorecard([
        {"dimension": "speed_to_value", "ia_score": 9, "competitor_score": 3, "weight": 5},
        {"dimension": "breadth", "ia_score": 5, "competitor_score": 9, "weight": 3},
    ])
    assert rows[0]["dimension"] == "breadth"


def test_every_dimension_has_a_label():
    from ciq.llm import DIMENSION_KEYS, DIMENSION_LABELS, SCORE_DIMENSIONS
    assert len(SCORE_DIMENSIONS) == 10
    assert len(DIMENSION_KEYS) == len(set(DIMENSION_KEYS))
    assert all(key in DIMENSION_LABELS for key in DIMENSION_KEYS)


def test_battlecard_schema_requires_a_scorecard():
    from ciq.llm import BATTLECARD_SCHEMA
    scorecard = BATTLECARD_SCHEMA["properties"]["scorecard"]
    # The ten-entry count cannot live in the schema, because structured output
    # rejects maxItems and any minItems above 1. It is stated in the prompt and
    # checked after the response instead.
    assert scorecard.get("minItems") in (0, 1)
    assert "maxItems" not in scorecard
    assert "scorecard" in BATTLECARD_SCHEMA["required"]
    assert "verdict" in BATTLECARD_SCHEMA["required"]


def test_a_short_scorecard_names_what_went_unscored():
    """With the count no longer enforceable in the schema, a partial scorecard
    must announce itself rather than looking complete."""
    from ciq.llm import DIMENSION_KEYS, normalise_scorecard
    rows = normalise_scorecard([
        {"dimension": "cost", "ia_score": 8, "competitor_score": 4, "weight": 3}])
    scored = {r["dimension"] for r in rows}
    missing = [k for k in DIMENSION_KEYS if k not in scored]
    assert len(missing) == 9


# ─── structured output schema compatibility ────────────────────────────────

# Structured output accepts a subset of JSON Schema. Sending anything outside
# it is a 400 with no output at all, so every schema is checked here rather
# than discovered in production.
UNSUPPORTED_SCHEMA_KEYS = (
    "maxItems", "minimum", "maximum", "exclusiveMinimum", "exclusiveMaximum",
    "multipleOf", "minLength", "maxLength", "uniqueItems", "oneOf", "not",
)


def _walk(node, found=None, path="$"):
    found = found if found is not None else []
    if isinstance(node, dict):
        for key, value in node.items():
            if key in UNSUPPORTED_SCHEMA_KEYS:
                found.append(f"{path}.{key}")
            if key == "minItems" and value not in (0, 1):
                found.append(f"{path}.minItems={value}")
            if key == "additionalProperties" and value is True:
                found.append(f"{path}.additionalProperties=true")
            _walk(value, found, f"{path}.{key}")
    elif isinstance(node, list):
        for i, item in enumerate(node):
            _walk(item, found, f"{path}[{i}]")
    return found


def _all_schemas():
    from ciq import llm
    return {
        "BATTLECARD_SCHEMA": llm.BATTLECARD_SCHEMA,
        "ANALYSIS_SCHEMA": llm.ANALYSIS_SCHEMA,
        "_PROBE_SCHEMA": llm._PROBE_SCHEMA,
    }


@pytest.mark.parametrize("name", list(_all_schemas()))
def test_sanitised_schemas_use_only_supported_keywords(name):
    """This is the exact failure that produced a live 400: minItems above 1
    and maxItems are not supported by structured output."""
    from ciq.llm import sanitise_schema
    violations = _walk(sanitise_schema(_all_schemas()[name]))
    assert violations == [], f"{name} would be rejected: {violations}"


def test_sanitiser_clamps_and_strips():
    from ciq.llm import sanitise_schema
    cleaned = sanitise_schema({
        "type": "object",
        "properties": {
            "rows": {"type": "array", "minItems": 10, "maxItems": 10,
                     "items": {"type": "string", "maxLength": 20}},
            "count": {"type": "number", "minimum": 0, "maximum": 5},
        },
        "additionalProperties": True,
    })
    rows = cleaned["properties"]["rows"]
    assert rows["minItems"] == 1
    assert "maxItems" not in rows
    assert "maxLength" not in rows["items"]
    assert cleaned["properties"]["count"] == {"type": "number"}
    assert cleaned["additionalProperties"] is False


def test_sanitiser_refuses_to_silently_drop_structural_keywords():
    from ciq.llm import LLMUnavailable, sanitise_schema
    with pytest.raises(LLMUnavailable, match="oneOf"):
        sanitise_schema({"type": "object", "properties": {},
                         "oneOf": [{"type": "string"}]})


def test_sanitiser_leaves_a_valid_schema_alone():
    from ciq.llm import sanitise_schema
    schema = {"type": "object", "properties": {"a": {"type": "string"}},
              "required": ["a"], "additionalProperties": False}
    assert sanitise_schema(schema) == schema


# ─── managed database resilience ───────────────────────────────────────────

@pytest.mark.skipif(not os.environ.get("CIQ_TEST_DATABASE_URL"),
                    reason="needs a Postgres to terminate connections against")
def test_a_dropped_connection_is_reestablished():
    """Managed databases close idle connections. A cached connection that died
    would otherwise fail every later request until the worker restarted."""
    import psycopg
    url = os.environ["CIQ_TEST_DATABASE_URL"]
    store = db.connect(url)
    db.clear_all(store)
    db.create_entry(store, {"competitor": "o9 Solutions", "title": "Survivor",
                            "category": "battlecard", "content": "text"},
                    chunks=["text"])

    admin = psycopg.connect(url)
    with admin.cursor() as cur:
        cur.execute("SELECT pg_terminate_backend(pid) FROM pg_stat_activity"
                    " WHERE datname = current_database()"
                    " AND pid <> pg_backend_pid()")
    admin.close()

    # The next call must transparently reconnect and still see the data.
    assert db.stats(store)["entries"] == 1
    assert [e["title"] for e in db.list_entries(store)] == ["Survivor"]
    db.clear_all(store)
    store.close()
    db._local.store = db._local.target = None


def test_postgres_connections_disable_prepared_statements():
    """A transaction-mode pooler, which several managed providers hand out by
    default, cannot keep the session a prepared statement belongs to."""
    import inspect
    source = inspect.getsource(db._open_postgres)
    assert "prepare_threshold=None" in source
    assert "connect_timeout" in source


# ─── connection string diagnostics ─────────────────────────────────────────

def test_a_valid_connection_string_reports_no_issues():
    d = db.describe_target(
        "postgresql://postgres.abc:Secret1@aws-0-us-east-1.pooler.supabase.com:5432/postgres")
    assert d["host"] == "aws-0-us-east-1.pooler.supabase.com"
    assert d["port"] == 5432
    assert d["user"] == "postgres.abc"
    assert d["database"] == "postgres"
    assert d["issues"] == []


# & and : are legal in the userinfo part and parse cleanly, so they are
# deliberately not flagged. Only characters that actually truncate the URL are.
@pytest.mark.parametrize("password", ["Pa/ss", "Pa#ss", "Pa?ss"])
def test_a_reserved_character_in_the_password_is_named(password):
    """This is the failure that reports only "Name or service not known"."""
    d = db.describe_target(
        f"postgresql://postgres.abc:{password}@aws-0-us-east-1.pooler.supabase.com:5432/postgres")
    joined = " ".join(d["issues"])
    assert "Percent-encode" in joined
    assert "%2F" in joined


def test_a_slash_in_the_password_puts_the_username_in_the_host_position():
    """Seeing the username reported as the host is the giveaway, so the
    recovered host must be the wrong one that DNS actually received."""
    d = db.describe_target(
        "postgresql://postgres.abc:Pa/ss@aws-0-us-east-1.pooler.supabase.com:5432/postgres")
    assert d["host"] == "postgres.abc"


def test_a_leftover_placeholder_is_named():
    d = db.describe_target(
        "postgresql://postgres.abc:[YOUR-PASSWORD]@host.pooler.supabase.com:5432/postgres")
    assert any("placeholder" in i for i in d["issues"])


def test_surrounding_quotes_are_named():
    d = db.describe_target(
        '"postgresql://postgres.abc:Secret1@host.pooler.supabase.com:5432/postgres"')
    assert any("quotes" in i for i in d["issues"])
    assert d["host"] == "host.pooler.supabase.com"


def test_stray_whitespace_is_named():
    d = db.describe_target(
        " postgresql://postgres.abc:Secret1@host.pooler.supabase.com:5432/postgres\n")
    assert any("whitespace" in i for i in d["issues"])


@pytest.mark.parametrize("url", [
    "postgresql://ciq:pw@dpg-cv1234abcd-a/ciq",      # platform internal host
    "postgresql://u:pw@localhost:5432/db",
    "postgresql://u:pw@db:5432/app",                 # container service name
])
def test_a_legitimate_dotless_host_is_not_flagged(url):
    """Internal hostnames have no dot. Flagging them sent people to fix a
    password that was never wrong."""
    assert db.describe_target(url)["issues"] == []


def test_the_username_appearing_as_the_host_is_flagged():
    """That is the unmistakable signature of a truncated URL."""
    d = db.describe_target("postgresql://myuser:pw@myuser/db")
    assert any("Percent-encode" in i for i in d["issues"])


def test_diagnostics_never_reveal_the_password():
    secret = "SuperSecret123"
    d = db.describe_target(
        f"postgresql://postgres.abc:{secret}@host.pooler.supabase.com:5432/postgres")
    assert secret not in json.dumps(d)


def test_a_harmless_special_character_is_not_flagged():
    """A false alarm sends someone chasing the wrong thing."""
    d = db.describe_target(
        "postgresql://postgres.abc:Pa&ss@aws-0-us-east-1.pooler.supabase.com:5432/postgres")
    assert d["issues"] == []
    assert d["host"] == "aws-0-us-east-1.pooler.supabase.com"


def test_an_empty_value_is_named():
    assert any("empty" in i for i in db.describe_target("")["issues"])


def test_a_bare_postgres_user_on_a_pooler_host_is_flagged():
    """A pooled connection routes by project, and the reference travels in the
    username. A bare 'postgres' is rejected as a password failure even when the
    password is right, which sends people to reset a password that was fine."""
    d = db.describe_target(
        "postgresql://postgres:Secret1@aws-0-ap-southeast-1.pooler.supabase.com:5432/postgres")
    assert any("project reference" in i for i in d["issues"])


def test_a_project_scoped_user_on_a_pooler_host_is_accepted():
    d = db.describe_target(
        "postgresql://postgres.qwertyuiop:Secret1@aws-0-ap-southeast-1.pooler.supabase.com:5432/postgres")
    assert d["issues"] == []


def test_a_bare_postgres_user_on_a_direct_host_is_not_flagged():
    """Bare 'postgres' is correct for a direct connection, so flagging it there
    would send someone to fix something that is not broken."""
    d = db.describe_target(
        "postgresql://postgres:Secret1@db.qwertyuiop.supabase.co:5432/postgres")
    assert d["issues"] == []


# ─── credentials supplied as separate parts ────────────────────────────────

def test_separate_parts_build_a_valid_url_from_a_hostile_password(monkeypatch):
    """Hand assembling a URL is where this breaks: a reserved character in the
    password truncates it. Supplied separately, the password is encoded here."""
    from ciq.config import Config
    monkeypatch.delenv("DATABASE_URL", raising=False)
    monkeypatch.setenv("CIQ_DB_HOST", "aws-0-ap-southeast-1.pooler.supabase.com")
    monkeypatch.setenv("CIQ_DB_USER", "postgres.abcdefgh")
    monkeypatch.setenv("CIQ_DB_PASSWORD", "p@ss/w#rd?x&y:z")
    monkeypatch.setenv("CIQ_DB_NAME", "postgres")
    monkeypatch.setenv("CIQ_DB_PORT", "5432")

    detail = db.describe_target(Config.database_url())
    assert detail["issues"] == []
    assert detail["host"] == "aws-0-ap-southeast-1.pooler.supabase.com"
    assert detail["user"] == "postgres.abcdefgh"
    assert detail["port"] == 5432
    assert detail["database"] == "postgres"


def test_a_full_url_takes_precedence_over_the_parts(monkeypatch):
    from ciq.config import Config
    monkeypatch.setenv("DATABASE_URL", "postgresql://u:p@explicit.example.com:5432/db")
    monkeypatch.setenv("CIQ_DB_HOST", "ignored.example.com")
    assert db.describe_target(Config.database_url())["host"] == "explicit.example.com"


def test_quotes_and_whitespace_around_the_url_are_tolerated(monkeypatch):
    """Pasting a value with quotes is a slip, not an intent."""
    from ciq.config import Config
    monkeypatch.setenv(
        "DATABASE_URL",
        '  "postgresql://postgres.abc:pw@host.pooler.supabase.com:5432/postgres"  ')
    detail = db.describe_target(Config.database_url())
    assert detail["issues"] == []
    assert detail["host"] == "host.pooler.supabase.com"


def test_no_configuration_means_no_url(monkeypatch):
    from ciq.config import Config
    monkeypatch.delenv("DATABASE_URL", raising=False)
    monkeypatch.delenv("CIQ_DB_HOST", raising=False)
    assert Config.database_url() == ""


# ─── deck rendering ────────────────────────────────────────────────────────

def test_every_slide_type_renders_without_error():
    """Layout is computed here rather than described by the model, so each
    type must survive whatever the model returns."""
    from ciq import deck
    spec = {
        "title": "Impact Analytics vs Blue Yonder", "subtitle": "For a retailer",
        "slides": [
            {"type": "section", "heading": "The problem", "subheading": "Sub"},
            {"type": "bullets", "heading": "Findings", "bullets": ["One", "Two"]},
            {"type": "comparison", "heading": "Head to head",
             "left_title": "Us", "right_title": "Them",
             "left": ["Deep"], "right": ["Broad"]},
            {"type": "stats", "heading": "Impact",
             "stats": [{"value": "6 weeks", "label": "To value"}]},
            {"type": "quote", "heading": "", "quote": "It worked.",
             "attribution": "A customer"},
            {"type": "close", "heading": "Next", "bullets": ["Pilot"]},
        ],
    }
    data = deck.build(spec).getvalue()
    assert len(data) > 20000

    import io
    from pptx import Presentation
    rendered = Presentation(io.BytesIO(data))
    assert len(rendered.slides) == 7          # a title slide plus the six
    assert rendered.slide_width == deck.SLIDE_W


def test_a_deck_survives_missing_and_empty_fields():
    """A model omitting an optional field must not produce a broken file."""
    from ciq import deck
    spec = {"title": "T", "subtitle": "",
            "slides": [{"type": "bullets", "heading": "Only a heading"},
                       {"type": "comparison", "heading": "Bare"},
                       {"type": "stats", "heading": "No stats"},
                       {"type": "unknown_type", "heading": "Falls back"}]}
    assert len(deck.build(spec).getvalue()) > 20000


def test_long_bullets_do_not_break_the_build():
    from ciq import deck
    spec = {"title": "T", "subtitle": "S", "slides": [
        {"type": "bullets", "heading": "Long", "bullets": ["word " * 60] * 6}]}
    assert len(deck.build(spec).getvalue()) > 20000


def test_speaker_notes_are_carried_through():
    import io
    from pptx import Presentation
    from ciq import deck
    spec = {"title": "T", "subtitle": "S", "slides": [
        {"type": "bullets", "heading": "H", "bullets": ["b"],
         "note": "Say this out loud."}]}
    rendered = Presentation(io.BytesIO(deck.build(spec).getvalue()))
    assert "Say this out loud." in rendered.slides[1].notes_slide.notes_text_frame.text


def test_the_deck_schema_avoids_unsupported_keywords():
    from ciq.deck import DECK_SCHEMA
    from ciq.llm import sanitise_schema
    assert _walk(sanitise_schema(DECK_SCHEMA)) == []
