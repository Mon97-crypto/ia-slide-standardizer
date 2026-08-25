"""Tests for the battlecard builder.

Run with: python3 -m pytest tests -q
"""

import os
import sys
import zipfile

import pytest

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from battlecards import brand, compat, library, schema, service  # noqa: E402
from battlecards.builder import build_presentation  # noqa: E402


@pytest.fixture
def deck_path(tmp_path):
    payload = library.scaffold('RELEX Solutions', 'InventorySmart')
    payload['proof_points'] = [{
        'stat': '18%', 'label': 'Forecast accuracy gain',
        'detail': 'A grocery chain lifted item level accuracy in one season.',
        'source': 'https://www.impactanalytics.co',
    }]
    result = service.build(payload, str(tmp_path))
    return result['path']


# ── copy rules ──────────────────────────────────────────────────────────────

def test_sanitize_removes_banned_dashes():
    assert schema.sanitize_text('Fast planning — every week') == 'Fast planning, every week'
    assert schema.sanitize_text('Range 10–20 units') == 'Range 10, 20 units'


def test_lint_flags_terminal_preposition():
    issues = schema.lint_copy('That is the gap they cannot deliver on.', 'test')
    assert any(issue['rule'] == 'terminal_preposition' for issue in issues)


def test_lint_flags_faq_heading():
    issues = schema.lint_copy('See the FAQ for detail.', 'test')
    assert any(issue['rule'] == 'faq_heading' for issue in issues)


def test_lint_flags_stale_statistic():
    issues = schema.lint_copy('Margins rose 12% in 2021.', 'test')
    assert any(issue['rule'] == 'stat_recency' for issue in issues)


def test_script_urls_are_dropped():
    card = schema.normalize({
        'meta': {'competitor': 'Acme'},
        'their_strengths': ['Broad suite'],
        'resources': [{'label': 'Bad', 'url': 'javascript:alert(1)'},
                      {'label': 'Good', 'url': 'www.impactanalytics.co'}],
    })
    assert card['resources'][0]['url'] == ''
    assert card['resources'][1]['url'] == 'https://www.impactanalytics.co'


def test_validate_requires_a_competitor_and_content():
    result = schema.validate(schema.normalize({}))
    assert len(result['errors']) == 2


# ── brand rules ─────────────────────────────────────────────────────────────

def test_one_solution_colour_per_card():
    for product, solution in brand.PRODUCT_SOLUTIONS.items():
        theme = brand.Theme.for_solution(brand.solution_for_product(product))
        assert theme.accent == brand.SOLUTION_COLORS[solution]
        assert theme.page_bg == brand.OFF_WHITE
        assert theme.ink == brand.BLACK


def test_readable_text_colour_on_brand_fills():
    assert brand.readable_on(brand.IMPACT_BLUE) == brand.WHITE
    assert brand.readable_on(brand.OFF_WHITE) == brand.BLACK
    assert brand.readable_on(brand.SOLUTION_COLORS['inventory_replenishment']) == brand.BLACK


def test_white_logo_variant_is_generated():
    path = brand.ensure_white_logo()
    assert path and os.path.exists(path)


def test_only_palette_colours_reach_the_deck(deck_path):
    allowed = {str(color) for color in [
        brand.IMPACT_BLUE, brand.OFF_WHITE, brand.BLACK, brand.WHITE,
        brand.ACCENT_ORANGE, brand.GRAY_1, brand.GRAY_2, brand.GRAY_3,
    ]}
    allowed |= {str(color) for color in brand.SOLUTION_COLORS.values()}
    # Tints of White over Impact Blue are the only derived values the builder uses.
    allowed |= {str(brand.mix(brand.WHITE, brand.IMPACT_BLUE, weight))
                for weight in (0.12, 0.35, 0.6)}

    import re
    used = set()
    with zipfile.ZipFile(deck_path) as archive:
        for name in archive.namelist():
            if name.startswith('ppt/slides/slide'):
                xml = archive.read(name).decode('utf-8')
                used |= set(re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', xml))
    assert used - allowed == set(), 'off palette colours: %s' % sorted(used - allowed)


# ── deck structure ──────────────────────────────────────────────────────────

def test_deck_builds_every_selected_section(deck_path):
    from pptx import Presentation
    presentation = Presentation(deck_path)
    assert len(presentation.slides._sldIdLst) >= 15


def test_slide_size_is_widescreen():
    card = schema.normalize(library.scaffold('Acme'))
    presentation = build_presentation(card)
    assert presentation.slide_width == 12192000
    assert presentation.slide_height == 6858000


def test_long_content_paginates():
    payload = library.scaffold('Acme', 'PriceSmart')
    payload['objections'] = [{'objection': 'Objection %d' % i, 'response': 'Answer', 'proof': 'Proof'}
                             for i in range(7)]
    card = schema.normalize(payload)
    presentation = build_presentation(card)
    titles = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == 'Objection handling':
                titles.append(shape.text_frame.text)
    assert len(titles) == 3  # seven objections spread over three slides


def test_sections_can_be_narrowed():
    payload = library.scaffold('Acme')
    payload['sections'] = ['one_pager']
    card = schema.normalize(payload)
    presentation = build_presentation(card)
    assert len(presentation.slides._sldIdLst) == 2  # cover is always included


def test_speaker_notes_can_be_switched_off():
    payload = library.scaffold('Acme')
    payload['options'] = {'include_notes': False}
    presentation = build_presentation(schema.normalize(payload))
    for slide in presentation.slides:
        if slide.has_notes_slide:
            assert slide.notes_slide.notes_text_frame.text.strip() == ''


# ── Google Slides compatibility ─────────────────────────────────────────────

def test_deck_passes_the_google_slides_audit(deck_path):
    report = compat.audit(deck_path)
    assert report['ok'], report['findings']
    assert report['findings'] == []


def test_every_run_names_a_typeface(deck_path):
    with zipfile.ZipFile(deck_path) as archive:
        for name in archive.namelist():
            if not name.startswith('ppt/slides/slide'):
                continue
            xml = archive.read(name).decode('utf-8')
            if '<a:t>' in xml:
                assert '<a:latin typeface="Inter Tight"' in xml, name


def test_no_autofit_font_scaling(deck_path):
    with zipfile.ZipFile(deck_path) as archive:
        for name in archive.namelist():
            if name.startswith('ppt/slides/slide'):
                assert 'fontScale' not in archive.read(name).decode('utf-8')


def test_tables_carry_no_theme_style(deck_path):
    with zipfile.ZipFile(deck_path) as archive:
        for name in archive.namelist():
            if name.startswith('ppt/slides/slide'):
                assert 'tableStyleId' not in archive.read(name).decode('utf-8')


def test_audit_reports_a_broken_deck(tmp_path):
    path = str(tmp_path / 'bad.pptx')
    from pptx import Presentation
    from pptx.util import Inches
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = 'No typeface here'
    presentation.save(path)
    report = compat.audit(path)
    assert any(item['rule'] == 'run_font' for item in report['findings'])


# ── service layer ───────────────────────────────────────────────────────────

def test_filename_is_safe(tmp_path):
    payload = library.scaffold('../../evil name', 'PriceSmart')
    result = service.build(payload, str(tmp_path))
    assert '/' not in result['filename'] and '..' not in result['filename']
    assert result['filename'].endswith('.pptx')


def test_build_rejects_an_empty_card(tmp_path):
    with pytest.raises(ValueError):
        service.build({}, str(tmp_path))


def test_presets_cover_every_solution():
    presets = service.presets()
    assert {entry['key'] for entry in presets['solutions']} == set(brand.SOLUTION_COLORS)
    assert len(presets['sections']) == len(schema.SECTION_ORDER)
